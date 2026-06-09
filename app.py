import os
import sqlite3
from urllib.parse import urlparse

try:
    import psycopg2
    import psycopg2.extras
except Exception:
    psycopg2 = None
from datetime import datetime, date, timedelta
from functools import wraps
from io import BytesIO
from email.message import EmailMessage
import smtplib
import base64
import requests
import re
import calendar

from roster_seed import OFFICIAL_ROSTER

from flask import Flask, render_template, request, redirect, url_for, flash, session, send_file, Response, jsonify
from werkzeug.security import generate_password_hash, check_password_hash
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None
    CategoryChartData = None
    XL_CHART_TYPE = None
    Inches = None
    Pt = None

try:
    from openai import OpenAI
except Exception:
    OpenAI = None

APP_TITLE = "Smart Late Students System"
DB_PATH = os.environ.get("DB_PATH", "late_students_dynamic.db")
DATABASE_URL = os.environ.get("DATABASE_URL", "").strip()
USE_POSTGRES = bool(DATABASE_URL)

app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "change-this-secret-key")

DEFAULT_ADMIN_USERNAME = os.environ.get("ADMIN_USERNAME", "admin")
DEFAULT_ADMIN_PASSWORD = os.environ.get("ADMIN_PASSWORD", "admin123")
UAE_TIME_OFFSET = timedelta(hours=4)


def now_uae():
    """Return current UAE time as a naive datetime for consistent display/storage."""
    return datetime.utcnow() + UAE_TIME_OFFSET


def uae_now_strings():
    now = now_uae()
    return now, now.strftime("%H:%M:%S"), now.isoformat(timespec="seconds")

OFFICIAL_GRADE_ORDER = {
    "Grade 9": 9,
    "Grade 10": 10,
    "Grade 11": 11,
    "Grade 12": 12,
}


def official_sections_by_grade():
    sections = {grade: set() for grade in OFFICIAL_GRADE_ORDER}
    for row in OFFICIAL_ROSTER:
        grade = row.get("grade_name")
        section = row.get("section_name")
        if grade and section:
            sections.setdefault(grade, set()).add(section)
    return {grade: sorted(vals, key=section_sort_key) for grade, vals in sections.items()}


def section_sort_key(name):
    """Sort sections like 9 ADV 1, 9 ADV 2, 9 GEN 1, 9 GEN 2."""
    text = str(name or "")
    grade_match = re.search(r"(\d+)", text)
    grade = int(grade_match.group(1)) if grade_match else 999
    program_order = 1 if "ADV" in text.upper() else 2 if "GEN" in text.upper() else 9
    nums = re.findall(r"\d+", text)
    number = int(nums[-1]) if nums else 999
    return (grade, program_order, number, text)



class PostgresConnection:
    def __init__(self, database_url):
        if psycopg2 is None:
            raise RuntimeError("psycopg2-binary is required for PostgreSQL. Add it to requirements.txt.")
        self.conn = psycopg2.connect(database_url, cursor_factory=psycopg2.extras.RealDictCursor)

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc, tb):
        if exc_type:
            self.conn.rollback()
        else:
            self.conn.commit()
        self.conn.close()

    def commit(self):
        self.conn.commit()

    def execute(self, sql, params=()):
        sql = adapt_sql_for_postgres(sql)
        cur = self.conn.cursor()
        cur.execute(sql, params or ())
        return cur


def adapt_sql_for_postgres(sql):
    """Translate the small SQLite SQL subset used by this app to PostgreSQL."""
    original = sql
    sql = sql.replace("?", "%s")
    sql = sql.replace("INTEGER PRIMARY KEY AUTOINCREMENT", "SERIAL PRIMARY KEY")
    sql = sql.replace("INSERT OR IGNORE INTO", "INSERT INTO")
    sql = sql.replace("INSERT OR REPLACE INTO", "INSERT INTO")

    compact = " ".join(sql.split())
    if "INSERT INTO grades" in compact and "ON CONFLICT" not in compact:
        sql += " ON CONFLICT (name) DO NOTHING"
    elif "INSERT INTO sections" in compact and "ON CONFLICT" not in compact:
        sql += " ON CONFLICT (grade_id, name) DO NOTHING"
    elif "INSERT INTO students" in compact and "ON CONFLICT" not in compact:
        sql += " ON CONFLICT (name, grade_id, section_id) DO NOTHING"
    elif "INSERT INTO late_records" in compact and "ON CONFLICT" not in compact:
        sql += " ON CONFLICT (student_id, late_date) DO UPDATE SET late_time=EXCLUDED.late_time, recorded_by=EXCLUDED.recorded_by, created_at=EXCLUDED.created_at"
    elif "INSERT INTO email_recipients" in compact and "ON CONFLICT" not in compact:
        sql += " ON CONFLICT (grade_id, email) DO UPDATE SET person_name=EXCLUDED.person_name, active=EXCLUDED.active, created_at=EXCLUDED.created_at"
    elif "INSERT INTO system_settings" in compact and "ON CONFLICT" not in compact:
        sql += " ON CONFLICT (key) DO UPDATE SET value=EXCLUDED.value"
    return sql


def get_conn():
    if USE_POSTGRES:
        return PostgresConnection(DATABASE_URL)
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    with get_conn() as conn:
        if not USE_POSTGRES:
            conn.execute("PRAGMA foreign_keys = ON")
        conn.execute("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT NOT NULL UNIQUE,
            password_hash TEXT NOT NULL,
            full_name TEXT NOT NULL,
            role TEXT NOT NULL CHECK(role IN ('system_owner','admin','user')),
            active INTEGER NOT NULL DEFAULT 1,
            created_at TEXT NOT NULL
        )
        """)
        conn.execute("""
        CREATE TABLE IF NOT EXISTS grades (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL UNIQUE,
            sort_order INTEGER NOT NULL DEFAULT 0
        )
        """)
        conn.execute("""
        CREATE TABLE IF NOT EXISTS sections (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            grade_id INTEGER NOT NULL,
            name TEXT NOT NULL,
            UNIQUE(grade_id, name),
            FOREIGN KEY(grade_id) REFERENCES grades(id) ON DELETE CASCADE
        )
        """)
        conn.execute("""
        CREATE TABLE IF NOT EXISTS students (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            grade_id INTEGER NOT NULL,
            section_id INTEGER NOT NULL,
            active INTEGER NOT NULL DEFAULT 1,
            created_at TEXT NOT NULL,
            UNIQUE(name, grade_id, section_id),
            FOREIGN KEY(grade_id) REFERENCES grades(id) ON DELETE CASCADE,
            FOREIGN KEY(section_id) REFERENCES sections(id) ON DELETE CASCADE
        )
        """)
        ensure_student_columns(conn)
        conn.execute("""
        CREATE TABLE IF NOT EXISTS late_records (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            student_id INTEGER NOT NULL,
            late_date TEXT NOT NULL,
            late_time TEXT NOT NULL,
            recorded_by INTEGER,
            created_at TEXT NOT NULL,
            UNIQUE(student_id, late_date),
            FOREIGN KEY(student_id) REFERENCES students(id) ON DELETE CASCADE,
            FOREIGN KEY(recorded_by) REFERENCES users(id) ON DELETE SET NULL
        )
        """)
        conn.execute("""
        CREATE TABLE IF NOT EXISTS email_recipients (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            grade_id INTEGER NOT NULL,
            person_name TEXT NOT NULL,
            email TEXT NOT NULL,
            active INTEGER NOT NULL DEFAULT 1,
            created_at TEXT NOT NULL,
            UNIQUE(grade_id, email),
            FOREIGN KEY(grade_id) REFERENCES grades(id) ON DELETE CASCADE
        )
        """)
        conn.execute("""
        CREATE TABLE IF NOT EXISTS system_settings (
            key TEXT PRIMARY KEY,
            value TEXT
        )
        """)
        ensure_role_schema(conn)
        ensure_seed_data(conn)
        ensure_single_system_owner(conn)
        ensure_default_settings(conn)
        conn.commit()



def table_columns(conn, table_name):
    """Return column names for SQLite or PostgreSQL."""
    if USE_POSTGRES:
        rows = conn.execute("""
            SELECT column_name AS name
            FROM information_schema.columns
            WHERE table_name=?
        """, (table_name,)).fetchall()
        return {r["name"] for r in rows}
    rows = conn.execute(f"PRAGMA table_info({table_name})").fetchall()
    return {r["name"] for r in rows}


def ensure_student_columns(conn):
    """Add official roster metadata columns to old databases without deleting data."""
    existing = table_columns(conn, "students")
    additions = [
        ("student_number", "TEXT"),
        ("name_ar", "TEXT"),
        ("name_en", "TEXT"),
        ("source_section", "TEXT"),
        ("source_grade", "TEXT"),
    ]
    for col, col_type in additions:
        if col not in existing:
            if USE_POSTGRES:
                conn.execute(f"ALTER TABLE students ADD COLUMN IF NOT EXISTS {col} {col_type}")
            else:
                conn.execute(f"ALTER TABLE students ADD COLUMN {col} {col_type}")


def upsert_official_student(conn, row, now):
    """Insert or update one official Excel student into the correct grade and section."""
    grade = conn.execute("SELECT id FROM grades WHERE name=?", (row["grade_name"],)).fetchone()
    if not grade:
        return
    section = conn.execute(
        "SELECT id FROM sections WHERE grade_id=? AND name=?",
        (grade["id"], row["section_name"]),
    ).fetchone()
    if not section:
        return
    display_name = row.get("name_en") or row.get("name_ar") or row.get("student_number")
    existing = None
    if row.get("student_number"):
        existing = conn.execute(
            "SELECT id FROM students WHERE student_number=?",
            (row["student_number"],),
        ).fetchone()
    if not existing:
        existing = conn.execute(
            "SELECT id FROM students WHERE name=? AND grade_id=? AND section_id=?",
            (display_name, grade["id"], section["id"]),
        ).fetchone()
    if existing:
        conn.execute(
            """
            UPDATE students
            SET name=?, student_number=?, name_ar=?, name_en=?, grade_id=?, section_id=?,
                source_section=?, source_grade=?, active=1
            WHERE id=?
            """,
            (
                display_name,
                row.get("student_number"),
                row.get("name_ar"),
                row.get("name_en"),
                grade["id"],
                section["id"],
                row.get("source_section"),
                row.get("source_grade"),
                existing["id"],
            ),
        )
    else:
        conn.execute(
            """
            INSERT INTO students(name, student_number, name_ar, name_en, grade_id, section_id, active, created_at, source_section, source_grade)
            VALUES (?, ?, ?, ?, ?, ?, 1, ?, ?, ?)
            """,
            (
                display_name,
                row.get("student_number"),
                row.get("name_ar"),
                row.get("name_en"),
                grade["id"],
                section["id"],
                now,
                row.get("source_section"),
                row.get("source_grade"),
            ),
        )


def ensure_role_schema(conn):
    """Allow the new system_owner role on PostgreSQL databases created by older versions."""
    if USE_POSTGRES:
        try:
            conn.execute("ALTER TABLE users DROP CONSTRAINT IF EXISTS users_role_check")
        except Exception:
            pass


def ensure_single_system_owner(conn):
    """Promote the default admin account to the single System Owner and prevent duplicates."""
    default_user = conn.execute("SELECT id FROM users WHERE username=?", (DEFAULT_ADMIN_USERNAME,)).fetchone()
    if default_user:
        try:
            conn.execute("UPDATE users SET role='system_owner', active=1 WHERE id=?", (default_user["id"],))
            conn.execute("UPDATE users SET role='admin' WHERE role='system_owner' AND id<>?", (default_user["id"],))
        except Exception:
            # Existing SQLite databases with an old CHECK constraint may not allow role migration.
            # Fresh deployments and PostgreSQL deployments are handled correctly.
            pass


def ensure_default_settings(conn):
    defaults = {
        "footer_owner_name": "",
        "footer_owner_section": "",
        "footer_extra": "Smart Late Students System",
    }
    for key, value in defaults.items():
        existing = conn.execute("SELECT key FROM system_settings WHERE key=?", (key,)).fetchone()
        if not existing:
            conn.execute("INSERT INTO system_settings(key, value) VALUES (?, ?)", (key, value))


def get_setting(key, default=""):
    with get_conn() as conn:
        row = conn.execute("SELECT value FROM system_settings WHERE key=?", (key,)).fetchone()
    return row["value"] if row and row["value"] is not None else default


def get_footer_settings():
    return {
        "owner_name": get_setting("footer_owner_name", ""),
        "owner_section": get_setting("footer_owner_section", ""),
        "extra": get_setting("footer_extra", "Smart Late Students System"),
    }


def is_admin_like(user):
    return bool(user and user["role"] in ("admin", "system_owner"))


def is_system_owner(user):
    return bool(user and user["role"] == "system_owner")

def ensure_seed_data(conn):
    now = datetime.now().isoformat(timespec="seconds")
    if not conn.execute("SELECT id FROM users WHERE username=?", (DEFAULT_ADMIN_USERNAME,)).fetchone():
        conn.execute(
            "INSERT INTO users(username,password_hash,full_name,role,active,created_at) VALUES (?,?,?,?,?,?)",
            (DEFAULT_ADMIN_USERNAME, generate_password_hash(DEFAULT_ADMIN_PASSWORD), "System Owner", "system_owner", 1, now),
        )

    # Official grade structure from the uploaded IDH Excel roster.
    for grade_name, order in OFFICIAL_GRADE_ORDER.items():
        conn.execute("INSERT OR IGNORE INTO grades(name, sort_order) VALUES (?, ?)", (grade_name, order))

    for grade_name, sections in official_sections_by_grade().items():
        grade = conn.execute("SELECT id FROM grades WHERE name=?", (grade_name,)).fetchone()
        if not grade:
            continue
        for sec in sections:
            conn.execute("INSERT OR IGNORE INTO sections(grade_id,name) VALUES (?,?)", (grade["id"], sec))

    # Import every student into the exact section from the Excel column, e.g.
    # 09[Adv-3rdLanguage]/1 -> Grade 9 / 9 ADV 1.
    for row in OFFICIAL_ROSTER:
        upsert_official_student(conn, row, now)



def current_user():
    uid = session.get("user_id")
    if not uid:
        return None
    with get_conn() as conn:
        return conn.execute("SELECT * FROM users WHERE id=? AND active=1", (uid,)).fetchone()


def login_required(fn):
    @wraps(fn)
    def wrapper(*args, **kwargs):
        if not current_user():
            return redirect(url_for("login"))
        return fn(*args, **kwargs)
    return wrapper


def admin_required(fn):
    @wraps(fn)
    def wrapper(*args, **kwargs):
        u = current_user()
        if not is_admin_like(u):
            flash("Admin access required.", "error")
            return redirect(url_for("dashboard"))
        return fn(*args, **kwargs)
    return wrapper


def system_owner_required(fn):
    @wraps(fn)
    def wrapper(*args, **kwargs):
        u = current_user()
        if not is_system_owner(u):
            flash("System Owner access required.", "error")
            return redirect(url_for("admin_home"))
        return fn(*args, **kwargs)
    return wrapper


@app.context_processor
def inject_user():
    return {"app_title": APP_TITLE, "me": current_user(), "footer_settings": get_footer_settings(), "is_admin_like": is_admin_like, "is_system_owner": is_system_owner}


@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "POST":
        username = request.form.get("username", "").strip()
        password = request.form.get("password", "")
        with get_conn() as conn:
            u = conn.execute("SELECT * FROM users WHERE username=? AND active=1", (username,)).fetchone()
        if u and check_password_hash(u["password_hash"], password):
            session.clear()
            session["user_id"] = u["id"]
            flash("Logged in successfully.", "success")
            return redirect(url_for("dashboard"))
        flash("Invalid username or password.", "error")
    return render_template("login.html")


@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))




@app.route("/account", methods=["GET", "POST"])
@login_required
def account_settings():
    me = current_user()
    if request.method == "POST":
        full_name = request.form.get("full_name", "").strip()
        username = request.form.get("username", "").strip()
        current_password = request.form.get("current_password", "")
        new_password = request.form.get("new_password", "")
        confirm_password = request.form.get("confirm_password", "")
        if not full_name or not username:
            flash("Full name and username are required.", "error")
            return redirect(url_for("account_settings"))
        if not check_password_hash(me["password_hash"], current_password):
            flash("Current password is incorrect.", "error")
            return redirect(url_for("account_settings"))
        if new_password and new_password != confirm_password:
            flash("New password and confirmation do not match.", "error")
            return redirect(url_for("account_settings"))
        try:
            with get_conn() as conn:
                existing = conn.execute("SELECT id FROM users WHERE username=? AND id<>?", (username, me["id"])).fetchone()
                if existing:
                    flash("This username is already used by another account.", "error")
                    return redirect(url_for("account_settings"))
                if new_password:
                    conn.execute("UPDATE users SET full_name=?, username=?, password_hash=? WHERE id=?", (full_name, username, generate_password_hash(new_password), me["id"]))
                else:
                    conn.execute("UPDATE users SET full_name=?, username=? WHERE id=?", (full_name, username, me["id"]))
                conn.commit()
            flash("Your login details were updated successfully.", "success")
            return redirect(url_for("account_settings"))
        except Exception as exc:
            flash(f"Could not update account: {exc}", "error")
            return redirect(url_for("account_settings"))
    return render_template("account.html")

@app.route("/")
@login_required
def dashboard():
    selected_grade = request.args.get("grade_id", type=int)
    selected_section = request.args.get("section_id", type=int)
    today_text = parse_date_or_today(request.args.get("date"))
    with get_conn() as conn:
        grades = conn.execute("SELECT * FROM grades ORDER BY sort_order, name").fetchall()
        if not selected_grade and grades:
            selected_grade = grades[0]["id"]
        sections = conn.execute("SELECT * FROM sections WHERE grade_id=? ORDER BY name", (selected_grade,)).fetchall() if selected_grade else []
        if not selected_section and sections:
            selected_section = sections[0]["id"]
        students = conn.execute("""
            SELECT st.*, g.name AS grade_name, s.name AS section_name
            FROM students st JOIN grades g ON g.id=st.grade_id JOIN sections s ON s.id=st.section_id
            WHERE st.active=1 AND st.grade_id=? AND st.section_id=?
            ORDER BY st.name
        """, (selected_grade, selected_section)).fetchall() if selected_grade and selected_section else []
        today_records = conn.execute("SELECT student_id, late_time FROM late_records WHERE late_date=?", (today_text,)).fetchall()
    record_map = {r["student_id"]: r["late_time"] for r in today_records}
    data = []
    for s in students:
        streak = warning_streak(s["id"], today_text)
        data.append({**dict(s), "late_time": record_map.get(s["id"]), "streak": streak, "color": status_color(streak)})
    return render_template("dashboard.html", grades=grades, sections=sections, students=data, selected_grade=selected_grade, selected_section=selected_section, today=today_text)


def consecutive_late_days(student_id, end_day):
    with get_conn() as conn:
        rows = conn.execute("SELECT late_date FROM late_records WHERE student_id=? AND late_date<=? ORDER BY late_date DESC", (student_id, end_day)).fetchall()
    days = {r["late_date"] for r in rows}
    d = datetime.fromisoformat(end_day).date()
    count = 0
    while d.isoformat() in days:
        count += 1
        d -= timedelta(days=1)
    return count


def parse_date_or_today(value):
    """Safely read yyyy-mm-dd from forms/query strings."""
    try:
        return datetime.strptime(value or "", "%Y-%m-%d").date().isoformat()
    except ValueError:
        return date.today().isoformat()


def total_late_days(student_id, end_day=None):
    """Total unique late days for a student up to the selected date.
    This is used for color status and the total shown on today's dashboard.
    """
    with get_conn() as conn:
        if end_day:
            row = conn.execute(
                "SELECT COUNT(DISTINCT late_date) AS c FROM late_records WHERE student_id=? AND late_date<=?",
                (student_id, end_day),
            ).fetchone()
        else:
            row = conn.execute(
                "SELECT COUNT(DISTINCT late_date) AS c FROM late_records WHERE student_id=?",
                (student_id,),
            ).fetchone()
    return row["c"] or 0


def warning_streak(student_id, today_text):
    # Backward-compatible name: now returns TOTAL late days up to selected date, not only consecutive days.
    return total_late_days(student_id, today_text)


def status_color(count):
    if count >= 4:
        return "red"
    if count >= 1:
        return "yellow"
    return "normal"


@app.post("/mark_late/<int:student_id>")
@login_required
def mark_late(student_id):
    selected_date = parse_date_or_today(request.form.get("date"))
    selected_grade = request.form.get("grade_id", type=int)
    selected_section = request.form.get("section_id", type=int)
    now, late_time_uae, created_at_uae = uae_now_strings()
    with get_conn() as conn:
        s = conn.execute("SELECT grade_id, section_id FROM students WHERE id=?", (student_id,)).fetchone()
        # The late_date can be a previous calendar date, while late_time records the real entry time.
        conn.execute("""
            INSERT OR REPLACE INTO late_records(student_id, late_date, late_time, recorded_by, created_at)
            VALUES (?, ?, ?, ?, ?)
        """, (student_id, selected_date, late_time_uae, session["user_id"], created_at_uae))
        conn.commit()
    flash(f"Late record saved for {selected_date}.", "success")
    return redirect(url_for("dashboard", grade_id=selected_grade or s["grade_id"], section_id=selected_section or s["section_id"], date=selected_date))


@app.post("/unmark_late/<int:student_id>")
@login_required
def unmark_late(student_id):
    selected_date = parse_date_or_today(request.form.get("date"))
    selected_grade = request.form.get("grade_id", type=int)
    selected_section = request.form.get("section_id", type=int)
    with get_conn() as conn:
        s = conn.execute("SELECT grade_id, section_id FROM students WHERE id=?", (student_id,)).fetchone()
        conn.execute("DELETE FROM late_records WHERE student_id=? AND late_date=?", (student_id, selected_date))
        conn.commit()
    flash(f"Late record removed for {selected_date}.", "success")
    return redirect(url_for("dashboard", grade_id=selected_grade or s["grade_id"], section_id=selected_section or s["section_id"], date=selected_date))


@app.route("/admin")
@login_required
@admin_required
def admin_home():
    page = max(request.args.get("page", 1, type=int), 1)
    per_page = request.args.get("per_page", 50, type=int)
    if per_page not in (30, 50, 100):
        per_page = 50
    offset = (page - 1) * per_page
    with get_conn() as conn:
        grades = conn.execute("SELECT * FROM grades ORDER BY sort_order, name").fetchall()
        sections = conn.execute("SELECT s.*, g.name AS grade_name FROM sections s JOIN grades g ON g.id=s.grade_id ORDER BY g.sort_order, s.name").fetchall()
        total_students_row = conn.execute("SELECT COUNT(*) AS c FROM students").fetchone()
        total_students = total_students_row["c"] or 0
        students = conn.execute("""
            SELECT st.*, g.name AS grade_name, s.name AS section_name
            FROM students st JOIN grades g ON g.id=st.grade_id JOIN sections s ON s.id=st.section_id
            ORDER BY st.active DESC, g.sort_order, s.name, st.name
            LIMIT ? OFFSET ?
        """, (per_page, offset)).fetchall()
        users = conn.execute("SELECT id, username, full_name, role, active, created_at FROM users ORDER BY role, username").fetchall()
    total_pages = max((total_students + per_page - 1) // per_page, 1)
    return render_template("admin.html", grades=grades, sections=sections, students=students, users=users, page=page, per_page=per_page, total_pages=total_pages, total_students=total_students, footer_settings=get_footer_settings())


@app.post("/admin/grades")
@login_required
@admin_required
def add_grade():
    with get_conn() as conn:
        conn.execute("INSERT OR IGNORE INTO grades(name, sort_order) VALUES (?, ?)", (request.form["name"].strip(), request.form.get("sort_order", type=int) or 0))
        conn.commit()
    return redirect(url_for("admin_home"))


@app.post("/admin/sections")
@login_required
@admin_required
def add_section():
    with get_conn() as conn:
        conn.execute("INSERT OR IGNORE INTO sections(grade_id, name) VALUES (?, ?)", (request.form.get("grade_id", type=int), request.form["name"].strip()))
        conn.commit()
    return redirect(url_for("admin_home"))


@app.post("/admin/students")
@login_required
@admin_required
def add_student():
    now = datetime.now().isoformat(timespec="seconds")
    section_id = request.form.get("section_id", type=int)
    with get_conn() as conn:
        sec = conn.execute("SELECT grade_id FROM sections WHERE id=?", (section_id,)).fetchone()
        name = request.form["name"].strip()
        conn.execute("""
            INSERT OR IGNORE INTO students(name, name_en, grade_id, section_id, active, created_at)
            VALUES (?, ?, ?, ?, 1, ?)
        """, (name, name, sec["grade_id"], section_id, now))
        conn.commit()
    return redirect(url_for("admin_home"))


@app.post("/admin/students/<int:student_id>/delete")
@login_required
@admin_required
def delete_student(student_id):
    with get_conn() as conn:
        conn.execute("UPDATE students SET active=0 WHERE id=?", (student_id,))
        conn.commit()
    return redirect(url_for("admin_home"))


@app.post("/admin/users")
@login_required
@admin_required
def add_user():
    me = current_user()
    requested_role = request.form.get("role", "user").strip()
    if requested_role == "system_owner":
        flash("Only one System Owner is allowed. This role cannot be created from the interface.", "error")
        return redirect(url_for("admin_home"))
    if me["role"] == "admin" and requested_role != "user":
        flash("Admins can add normal users only. System Owner controls admin-level accounts.", "error")
        return redirect(url_for("admin_home"))
    now = now_uae().isoformat(timespec="seconds")
    with get_conn() as conn:
        conn.execute("INSERT INTO users(username,password_hash,full_name,role,active,created_at) VALUES (?,?,?,?,1,?)", (
            request.form["username"].strip(), generate_password_hash(request.form["password"]), request.form["full_name"].strip(), requested_role, now
        ))
        conn.commit()
    return redirect(url_for("admin_home"))


@app.post("/admin/users/<int:user_id>/toggle")
@login_required
@system_owner_required
def toggle_user(user_id):
    if user_id == session.get("user_id"):
        flash("You cannot deactivate your own account.", "error")
        return redirect(url_for("admin_home"))
    with get_conn() as conn:
        target = conn.execute("SELECT role FROM users WHERE id=?", (user_id,)).fetchone()
        if target and target["role"] == "system_owner":
            flash("The System Owner account cannot be deactivated.", "error")
            return redirect(url_for("admin_home"))
        conn.execute("UPDATE users SET active = CASE active WHEN 1 THEN 0 ELSE 1 END WHERE id=?", (user_id,))
        conn.commit()
    return redirect(url_for("admin_home"))


@app.post("/owner/footer")
@login_required
@system_owner_required
def update_footer_settings():
    values = {
        "footer_owner_name": request.form.get("owner_name", "").strip(),
        "footer_owner_section": request.form.get("owner_section", "").strip(),
        "footer_extra": request.form.get("extra", "").strip(),
    }
    with get_conn() as conn:
        for key, value in values.items():
            conn.execute("INSERT OR REPLACE INTO system_settings(key, value) VALUES (?, ?)", (key, value))
        conn.commit()
    flash("Footer copyright details updated.", "success")
    return redirect(url_for("admin_home"))




@app.post("/owner/users/<int:user_id>/update")
@login_required
@system_owner_required
def owner_update_user(user_id):
    full_name = request.form.get("full_name", "").strip()
    username = request.form.get("username", "").strip()
    requested_role = request.form.get("role", "user").strip()
    new_password = request.form.get("new_password", "")
    next_url = request.form.get("next") or url_for("admin_home")
    if not full_name or not username:
        flash("Full name and username are required.", "error")
        return redirect(next_url)
    if requested_role not in ("user", "admin", "system_owner"):
        flash("Invalid role selected.", "error")
        return redirect(next_url)
    with get_conn() as conn:
        target = conn.execute("SELECT * FROM users WHERE id=?", (user_id,)).fetchone()
        if not target:
            flash("User not found.", "error")
            return redirect(next_url)
        # Only the existing owner account may remain system_owner. No account can be promoted to owner.
        if target["role"] == "system_owner":
            final_role = "system_owner"
        elif requested_role == "system_owner":
            flash("Only one System Owner is allowed. You can promote users to Admin or downgrade them to User only.", "error")
            return redirect(next_url)
        else:
            final_role = requested_role
        existing = conn.execute("SELECT id FROM users WHERE username=? AND id<>?", (username, user_id)).fetchone()
        if existing:
            flash("This username is already used by another account.", "error")
            return redirect(next_url)
        if new_password:
            conn.execute("UPDATE users SET full_name=?, username=?, role=?, password_hash=? WHERE id=?", (full_name, username, final_role, generate_password_hash(new_password), user_id))
        else:
            conn.execute("UPDATE users SET full_name=?, username=?, role=? WHERE id=?", (full_name, username, final_role, user_id))
        conn.commit()
    flash("Login details updated by System Owner.", "success")
    return redirect(next_url)



@app.post("/owner/users/<int:user_id>/delete")
@login_required
@system_owner_required
def owner_delete_user(user_id):
    if user_id == session.get("user_id"):
        flash("You cannot delete your own System Owner account.", "error")
        return redirect(url_for("admin_home"))
    with get_conn() as conn:
        target = conn.execute("SELECT role FROM users WHERE id=?", (user_id,)).fetchone()
        if not target:
            flash("User not found.", "error")
            return redirect(url_for("admin_home"))
        if target["role"] == "system_owner":
            flash("The System Owner account is protected and cannot be deleted.", "error")
            return redirect(url_for("admin_home"))
        conn.execute("DELETE FROM users WHERE id=?", (user_id,))
        conn.commit()
    flash("User deleted by System Owner.", "success")
    return redirect(url_for("admin_home"))


@app.post("/owner/grades/<int:grade_id>/delete")
@login_required
@system_owner_required
def owner_delete_grade(grade_id):
    with get_conn() as conn:
        grade = conn.execute("SELECT name FROM grades WHERE id=?", (grade_id,)).fetchone()
        if not grade:
            flash("Grade not found.", "error")
            return redirect(url_for("admin_home"))
        conn.execute("DELETE FROM late_records WHERE student_id IN (SELECT id FROM students WHERE grade_id=?)", (grade_id,))
        conn.execute("DELETE FROM students WHERE grade_id=?", (grade_id,))
        conn.execute("DELETE FROM sections WHERE grade_id=?", (grade_id,))
        conn.execute("DELETE FROM email_recipients WHERE grade_id=?", (grade_id,))
        conn.execute("DELETE FROM grades WHERE id=?", (grade_id,))
        conn.commit()
    flash(f"Grade deleted: {grade['name']}. Related sections, students, late records, and recipients were removed.", "success")
    return redirect(url_for("admin_home"))


@app.post("/owner/sections/<int:section_id>/delete")
@login_required
@system_owner_required
def owner_delete_section(section_id):
    with get_conn() as conn:
        section = conn.execute("SELECT s.name, g.name AS grade_name FROM sections s JOIN grades g ON g.id=s.grade_id WHERE s.id=?", (section_id,)).fetchone()
        if not section:
            flash("Section not found.", "error")
            return redirect(url_for("admin_home"))
        conn.execute("DELETE FROM late_records WHERE student_id IN (SELECT id FROM students WHERE section_id=?)", (section_id,))
        conn.execute("DELETE FROM students WHERE section_id=?", (section_id,))
        conn.execute("DELETE FROM sections WHERE id=?", (section_id,))
        conn.commit()
    flash(f"Section deleted: {section['grade_name']} - {section['name']}. Related students and late records were removed.", "success")
    return redirect(url_for("admin_home"))


@app.post("/admin/email-recipients")
@login_required
@admin_required
def add_email_recipient():
    now = datetime.now().isoformat(timespec="seconds")
    grade_id = request.form.get("grade_id", type=int)
    person_name = request.form.get("person_name", "").strip()
    email = request.form.get("email", "").strip()
    if not grade_id or not person_name or not email:
        flash("Grade, recipient name, and email are required.", "error")
        return redirect(url_for("admin_home"))
    with get_conn() as conn:
        conn.execute("""
            INSERT OR REPLACE INTO email_recipients(grade_id, person_name, email, active, created_at)
            VALUES (?, ?, ?, 1, ?)
        """, (grade_id, person_name, email, now))
        conn.commit()
    flash("Email recipient saved successfully.", "success")
    return redirect(url_for("admin_email"))


@app.post("/admin/email-recipients/<int:recipient_id>/toggle")
@login_required
@admin_required
def toggle_email_recipient(recipient_id):
    with get_conn() as conn:
        conn.execute("UPDATE email_recipients SET active = CASE active WHEN 1 THEN 0 ELSE 1 END WHERE id=?", (recipient_id,))
        conn.commit()
    return redirect(url_for("admin_email"))


@app.route("/admin/email", methods=["GET", "POST"])
@login_required
@admin_required
def admin_email():
    today_iso = date.today().isoformat()
    with get_conn() as conn:
        grades = conn.execute("SELECT * FROM grades ORDER BY sort_order, name").fetchall()
        recipients = conn.execute("""
            SELECT er.*, g.name AS grade_name
            FROM email_recipients er JOIN grades g ON g.id=er.grade_id
            ORDER BY g.sort_order, er.person_name
        """).fetchall()
    if request.method == "POST":
        grade_id = request.form.get("grade_id", type=int)
        recipient_id = request.form.get("recipient_id", type=int)
        day_text = parse_date_or_today(request.form.get("report_date") or today_iso)
        ok, message = send_grade_email(grade_id, recipient_id, day_text)
        flash(message, "success" if ok else "error")
        return redirect(url_for("admin_email"))
    return render_template("admin_email.html", grades=grades, recipients=recipients, today=today_iso)


@app.route("/report")
@login_required
def report():
    today_iso = date.today().isoformat()
    from_date = parse_date_or_today(request.args.get("from_date") or request.args.get("date") or today_iso)
    to_date = parse_date_or_today(request.args.get("to_date") or from_date)
    if from_date > to_date:
        from_date, to_date = to_date, from_date
    page = max(request.args.get("page", 1, type=int), 1)
    per_page = request.args.get("per_page", 50, type=int)
    if per_page not in (30, 50, 100):
        per_page = 50
    records_all = get_records_range(from_date, to_date)
    total_records = len(records_all)
    total_pages = max((total_records + per_page - 1) // per_page, 1)
    if page > total_pages:
        page = total_pages
    start = (page - 1) * per_page
    records = records_all[start:start + per_page]
    return render_template("report.html", records=records, from_date=from_date, to_date=to_date, today=today_iso, page=page, per_page=per_page, total_pages=total_pages, total_records=total_records, start_index=start)


@app.route("/student/<int:student_id>/late-details")
@login_required
def student_late_details(student_id):
    today_iso = date.today().isoformat()
    filter_key = request.args.get("filter", "all_time")
    from_arg = request.args.get("from_date")
    to_arg = request.args.get("to_date")
    from_date, to_date, label = resolve_student_detail_period(student_id, filter_key, from_arg, to_arg)
    with get_conn() as conn:
        student = conn.execute("""
            SELECT s.*, g.name AS grade_name, sec.name AS section_name
            FROM students s
            JOIN grades g ON g.id=s.grade_id
            JOIN sections sec ON sec.id=s.section_id
            WHERE s.id=?
        """, (student_id,)).fetchone()
        rows = conn.execute("""
            SELECT r.*, u.full_name AS recorder_name, u.username AS recorder_username
            FROM late_records r
            LEFT JOIN users u ON u.id=r.recorded_by
            WHERE r.student_id=? AND r.late_date BETWEEN ? AND ?
            ORDER BY r.late_date DESC, r.late_time DESC
        """, (student_id, from_date, to_date)).fetchall()
    if not student:
        flash("Student not found.", "danger")
        return redirect(url_for("report", from_date=from_date, to_date=to_date))
    records = format_student_detail_records(rows)
    return render_template("student_details.html", student=dict(student), records=records, from_date=from_date, to_date=to_date, period_label=label, total_days=len(records))


@app.route("/student/<int:student_id>/late-details-data")
@login_required
def student_late_details_data(student_id):
    filter_key = request.args.get("filter", "selected_period")
    from_arg = request.args.get("from_date")
    to_arg = request.args.get("to_date")
    from_date, to_date, label = resolve_student_detail_period(student_id, filter_key, from_arg, to_arg)
    with get_conn() as conn:
        student = conn.execute("""
            SELECT s.*, g.name AS grade_name, sec.name AS section_name
            FROM students s
            JOIN grades g ON g.id=s.grade_id
            JOIN sections sec ON sec.id=s.section_id
            WHERE s.id=?
        """, (student_id,)).fetchone()
        if not student:
            return jsonify({"ok": False, "message": "Student not found."}), 404
        rows = conn.execute("""
            SELECT r.*, u.full_name AS recorder_name, u.username AS recorder_username
            FROM late_records r
            LEFT JOIN users u ON u.id=r.recorded_by
            WHERE r.student_id=? AND r.late_date BETWEEN ? AND ?
            ORDER BY r.late_date DESC, r.late_time DESC
        """, (student_id, from_date, to_date)).fetchall()
    return jsonify({
        "ok": True,
        "student": {
            "id": student["id"],
            "student_number": student["student_number"] or "-",
            "name": student["name_en"] or student["name"] or "-",
            "arabic_name": student["name_ar"] or "-",
            "grade": student["grade_name"],
            "section": student["section_name"],
            "created_at": student["created_at"] or "-",
        },
        "period": {"from_date": from_date, "to_date": to_date, "label": label},
        "total_days": len(rows),
        "records": format_student_detail_records(rows),
    })


def format_student_detail_records(rows):
    records = []
    for r in rows:
        item = dict(r)
        item["day_name"] = datetime.strptime(item["late_date"], "%Y-%m-%d").strftime("%A")
        item["recorded_by_name"] = item.get("recorder_name") or "Unknown"
        item["recorded_by_username"] = item.get("recorder_username") or "-"
        records.append(item)
    return records


def resolve_student_detail_period(student_id, filter_key="selected_period", from_arg=None, to_arg=None):
    today_iso = date.today().isoformat()
    filter_key = (filter_key or "selected_period").strip()
    if filter_key == "last_week":
        end = date.today()
        start = end - timedelta(days=6)
        return start.isoformat(), end.isoformat(), "Last 7 days"
    if filter_key == "last_month":
        first_this_month = date.today().replace(day=1)
        last_prev_month = first_this_month - timedelta(days=1)
        first_prev_month = last_prev_month.replace(day=1)
        return first_prev_month.isoformat(), last_prev_month.isoformat(), "Last month"
    if filter_key == "all_time":
        with get_conn() as conn:
            student = conn.execute("SELECT created_at FROM students WHERE id=?", (student_id,)).fetchone()
            earliest = conn.execute("SELECT MIN(late_date) AS first_late FROM late_records WHERE student_id=?", (student_id,)).fetchone()
        created = (student["created_at"] or today_iso)[:10] if student else today_iso
        first_late = earliest["first_late"] if earliest and earliest["first_late"] else created
        start = min(created, first_late)
        return start, today_iso, "Whole period"
    # selected_period is the report period; never defaults to 1900.
    from_date = parse_date_or_today(from_arg or today_iso)
    to_date = parse_date_or_today(to_arg or today_iso)
    if from_date > to_date:
        from_date, to_date = to_date, from_date
    return from_date, to_date, "Selected report period"


def get_records_range(from_date, to_date, grade_id=None, section_id=None):
    params = [from_date, to_date]
    filters = []
    if grade_id:
        filters.append("g.id=?")
        params.append(int(grade_id))
    if section_id:
        filters.append("sec.id=?")
        params.append(int(section_id))
    extra_filter = (" AND " + " AND ".join(filters)) if filters else ""
    with get_conn() as conn:
        rows = conn.execute(f"""
            SELECT r.*, st.name AS student_name, st.student_number AS student_number, st.name_ar AS student_name_ar,
                   st.name_en AS student_name_en, g.id AS grade_id, g.name AS grade_name,
                   sec.id AS section_id, sec.name AS section_name,
                   u.full_name AS recorder_name, u.username AS recorder_username
            FROM late_records r
            JOIN students st ON st.id=r.student_id
            JOIN grades g ON g.id=st.grade_id
            JOIN sections sec ON sec.id=st.section_id
            LEFT JOIN users u ON u.id=r.recorded_by
            WHERE r.late_date BETWEEN ? AND ?{extra_filter}
            ORDER BY r.late_date, g.sort_order, sec.name, st.name, r.late_time
        """, params).fetchall()
    result = []
    for r in rows:
        item = dict(r)
        item["total_days"] = total_late_days(r["student_id"], to_date)
        item["day_name"] = datetime.strptime(r["late_date"], "%Y-%m-%d").strftime("%A")
        result.append(item)
    return result


def get_daily_records(day_text):
    return get_records_range(day_text, day_text)


@app.route("/export/late_report.xlsx")
@login_required
def export_late_report_excel():
    today_iso = date.today().isoformat()
    from_date = parse_date_or_today(request.args.get("from_date") or request.args.get("date") or today_iso)
    to_date = parse_date_or_today(request.args.get("to_date") or from_date)
    if from_date > to_date:
        from_date, to_date = to_date, from_date
    grade_id = request.args.get("grade_id", type=int)
    records = get_records_range(from_date, to_date, grade_id=grade_id)
    return build_excel_response(records, from_date, to_date, grade_id=grade_id)


@app.route("/export/daily.xlsx")
@login_required
def export_daily_excel():
    # Kept for old buttons/links; now supports date or from_date/to_date.
    return export_late_report_excel()


def sanitize_sheet_name(name):
    """Clean invalid Excel sheet-name characters and keep length <= 31."""
    bad_chars = ['\\', '/', '?', '*', '[', ']', ':']
    clean = str(name or "Sheet")
    for ch in bad_chars:
        clean = clean.replace(ch, '-')
    return clean[:31] or "Sheet"


def each_date_between(from_date, to_date):
    start = datetime.strptime(from_date, "%Y-%m-%d").date()
    end = datetime.strptime(to_date, "%Y-%m-%d").date()
    while start <= end:
        yield start.isoformat(), start.strftime("%A")
        start += timedelta(days=1)


def all_grades_for_excel(grade_id=None):
    with get_conn() as conn:
        if grade_id:
            return conn.execute("SELECT id, name, sort_order FROM grades WHERE id=? ORDER BY sort_order, name", (int(grade_id),)).fetchall()
        return conn.execute("SELECT id, name, sort_order FROM grades ORDER BY sort_order, name").fetchall()


def build_excel_workbook_bytes(records, from_date, to_date, grade_id=None):
    wb = Workbook()
    default_ws = wb.active
    wb.remove(default_ws)

    grades = all_grades_for_excel(grade_id)
    records_by_grade_date = {}
    for r in records:
        records_by_grade_date.setdefault(r["grade_name"], {}).setdefault(r["late_date"], []).append(r)

    # Professional colors
    dark_blue = "1F4E78"
    light_blue = "5B9BD5"
    pale_blue = "D9EAF7"
    yellow = "FFF2CC"
    red = "F4CCCC"
    white = "FFFFFF"
    grey = "F2F2F2"
    border_color = "B7B7B7"

    headers = [
        "No.",
        "Student No.",
        "English Name",
        "Arabic Name",
        "Grade",
        "Section",
        "Late/Absent Day",
        "Date",
        "Late Time",
        "Total Late Days",
        "Recorded By",
        "Username",
        "Entry Created At",
    ]

    if not grades:
        # Fallback sheet in the unlikely case the admin has removed all grades.
        grades = [{"name": "No Grades", "sort_order": 0}]

    used_sheet_names = set()
    for grade in grades:
        grade_name = grade["name"]
        base_name = sanitize_sheet_name(grade_name)
        sheet_name = base_name
        n = 2
        while sheet_name in used_sheet_names:
            suffix = f" {n}"
            sheet_name = sanitize_sheet_name(base_name[:31-len(suffix)] + suffix)
            n += 1
        used_sheet_names.add(sheet_name)

        ws = wb.create_sheet(sheet_name)
        ws.sheet_view.rightToLeft = False
        title = f"{grade_name} - Late/Absent Interval Report"
        period = f"Period: {from_date}" if from_date == to_date else f"Period: {from_date} to {to_date}"

        ws.merge_cells("A1:M1")
        ws["A1"] = title
        ws["A1"].font = Font(size=16, bold=True, color=white)
        ws["A1"].fill = PatternFill("solid", fgColor=dark_blue)
        ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
        ws.row_dimensions[1].height = 28

        ws.merge_cells("A2:M2")
        total_records_for_grade = sum(len(v) for v in records_by_grade_date.get(grade_name, {}).values())
        ws["A2"] = f"{period} | Total records in this grade: {total_records_for_grade}"
        ws["A2"].font = Font(size=11, bold=True, color="1F1F1F")
        ws["A2"].fill = PatternFill("solid", fgColor=pale_blue)
        ws["A2"].alignment = Alignment(horizontal="center", vertical="center")

        header_row = 4
        for col, header in enumerate(headers, start=1):
            cell = ws.cell(header_row, col)
            cell.value = header
            cell.font = Font(bold=True, color=white)
            cell.fill = PatternFill("solid", fgColor=dark_blue)
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        ws.row_dimensions[header_row].height = 24
        ws.freeze_panes = "A5"

        counter = 1
        records_for_grade = records_by_grade_date.get(grade_name, {})
        for day_iso, day_name in each_date_between(from_date, to_date):
            group_row = ws.max_row + 1
            ws.merge_cells(start_row=group_row, start_column=1, end_row=group_row, end_column=13)
            group_cell = ws.cell(group_row, 1)
            group_cell.value = f"{day_name} - {day_iso}"
            group_cell.font = Font(bold=True, color=white)
            group_cell.fill = PatternFill("solid", fgColor=light_blue)
            group_cell.alignment = Alignment(horizontal="left", vertical="center")
            ws.row_dimensions[group_row].height = 21

            day_records = records_for_grade.get(day_iso, [])
            if day_records:
                for r in day_records:
                    ws.append([
                        counter,
                        r.get("student_number") or "-",
                        r["student_name"],
                        r.get("student_name_ar") or "-",
                        r["grade_name"],
                        r["section_name"],
                        r["day_name"],
                        r["late_date"],
                        r["late_time"],
                        r["total_days"],
                        r["recorder_name"] or "Unknown",
                        r["recorder_username"] or "-",
                        r["created_at"],
                    ])
                    data_row = ws.max_row
                    total_cell = ws.cell(data_row, 10)
                    if (r["total_days"] or 0) >= 4:
                        total_cell.fill = PatternFill("solid", fgColor=red)
                    elif (r["total_days"] or 0) >= 1:
                        total_cell.fill = PatternFill("solid", fgColor=yellow)
                    counter += 1
            else:
                ws.append(["", "", "No late/absent students recorded", "", grade_name, "", day_name, day_iso, "", "", "", "", ""])
                no_row = ws.max_row
                for col in range(1, 14):
                    ws.cell(no_row, col).fill = PatternFill("solid", fgColor=grey)
                ws.cell(no_row, 2).font = Font(italic=True, color="666666")

        # Formatting and print readiness
        thin = Side(style="thin", color=border_color)
        for row in ws.iter_rows(min_row=4, max_row=ws.max_row, min_col=1, max_col=13):
            for cell in row:
                cell.border = Border(top=thin, bottom=thin, left=thin, right=thin)
                cell.alignment = Alignment(vertical="center", wrap_text=True)

        widths = [7, 15, 30, 28, 13, 14, 18, 14, 13, 17, 24, 16, 22]
        for i, width in enumerate(widths, start=1):
            ws.column_dimensions[get_column_letter(i)].width = width
        ws.auto_filter.ref = f"A4:M{ws.max_row}"
        ws.page_setup.orientation = "landscape"
        ws.page_setup.fitToWidth = 1
        ws.page_setup.fitToHeight = 0
        ws.sheet_properties.pageSetUpPr.fitToPage = True
        ws.print_title_rows = "1:4"
        ws.page_margins.left = 0.25
        ws.page_margins.right = 0.25
        ws.page_margins.top = 0.5
        ws.page_margins.bottom = 0.5

    bio = BytesIO()
    wb.save(bio)
    bio.seek(0)
    return bio


def build_excel_response(records, from_date, to_date, grade_id=None):
    bio = build_excel_workbook_bytes(records, from_date, to_date, grade_id=grade_id)
    suffix = from_date if from_date == to_date else f"{from_date}_to_{to_date}"
    grade_part = "all_grades"
    if grade_id:
        with get_conn() as conn:
            g = conn.execute("SELECT name FROM grades WHERE id=?", (int(grade_id),)).fetchone()
        grade_part = sanitize_sheet_name(g["name"] if g else f"grade_{grade_id}").replace(" ", "_")
    return send_file(
        bio,
        as_attachment=True,
        download_name=f"late_absent_{grade_part}_{suffix}.xlsx",
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )



# =========================
# AI ANALYTICS + CHATBOT
# =========================

def grade_number_from_name(name):
    digits = ''.join(ch for ch in str(name or '') if ch.isdigit())
    return int(digits) if digits else 999


def get_ai_risk_report(from_date, to_date):
    """Rule-based AI engine for lateness risk.

    Important upgrade:
    - The selected date range is used for the current report totals.
    - Repeated weekday pattern is checked historically up to the selected To Date,
      not only inside the selected range. This means if a student was late on a
      Friday this week and also on a Friday in any previous week, the AI will flag
      Friday as a repeated pattern even when the current range has only one Friday.
    """
    with get_conn() as conn:
        students = conn.execute("""
            SELECT st.id, st.name AS student_name, g.name AS grade_name, sec.name AS section_name
            FROM students st
            JOIN grades g ON g.id=st.grade_id
            JOIN sections sec ON sec.id=st.section_id
            WHERE st.active=1
            ORDER BY g.sort_order, sec.name, st.name
        """).fetchall()

        # Records inside the selected interval: these are the records shown in the report.
        rows = conn.execute("""
            SELECT r.student_id, r.late_date, r.late_time, u.full_name AS recorder_name, u.username AS recorder_username
            FROM late_records r
            LEFT JOIN users u ON u.id=r.recorded_by
            WHERE r.late_date BETWEEN ? AND ?
            ORDER BY r.student_id, r.late_date
        """, (from_date, to_date)).fetchall()

        # Historical records up to the selected To Date: used only for repeated weekday pattern detection.
        historical_rows = conn.execute("""
            SELECT r.student_id, r.late_date, r.late_time, u.full_name AS recorder_name, u.username AS recorder_username
            FROM late_records r
            LEFT JOIN users u ON u.id=r.recorded_by
            WHERE r.late_date <= ?
            ORDER BY r.student_id, r.late_date
        """, (to_date,)).fetchall()

    by_student = {}
    for r in rows:
        by_student.setdefault(r['student_id'], []).append(dict(r))

    historical_by_student = {}
    for r in historical_rows:
        historical_by_student.setdefault(r['student_id'], []).append(dict(r))

    report = []
    for st in students:
        records = by_student.get(st['id'], [])
        historical_records = historical_by_student.get(st['id'], [])

        dates = sorted({r['late_date'] for r in records})
        historical_dates = sorted({r['late_date'] for r in historical_records})
        total = len(dates)
        historical_total = len(historical_dates)

        # Pattern inside the selected interval.
        interval_weekday_counts = {}
        for d in dates:
            day_name = datetime.strptime(d, "%Y-%m-%d").strftime("%A")
            interval_weekday_counts[day_name] = interval_weekday_counts.get(day_name, 0) + 1
        repeated_weekdays = {k: v for k, v in interval_weekday_counts.items() if v >= 2}

        # Historical pattern across previous weeks up to To Date.
        # This is the requested logic: one Tuesday last week + one Tuesday this week = repeated pattern.
        historical_weekday_dates = {}
        for d in historical_dates:
            day_name = datetime.strptime(d, "%Y-%m-%d").strftime("%A")
            historical_weekday_dates.setdefault(day_name, []).append(d)
        historical_repeated_weekdays = {
            day: {
                "count": len(day_dates),
                "dates": day_dates,
                "latest_dates": day_dates[-5:],
            }
            for day, day_dates in historical_weekday_dates.items()
            if len(day_dates) >= 2
        }

        # Combined pattern shown to the user: historical is more useful, interval is still kept for compatibility.
        pattern_for_scoring = historical_repeated_weekdays or {k: {"count": v, "dates": [], "latest_dates": []} for k, v in repeated_weekdays.items()}

        max_consecutive = 0
        current = 0
        previous_day = None
        for d in dates:
            current_day = datetime.strptime(d, "%Y-%m-%d").date()
            if previous_day and (current_day - previous_day).days == 1:
                current += 1
            else:
                current = 1
            max_consecutive = max(max_consecutive, current)
            previous_day = current_day

        # Historical consecutive streak is also useful if the selected range is small.
        historical_max_consecutive = 0
        historical_current = 0
        historical_previous_day = None
        for d in historical_dates:
            current_day = datetime.strptime(d, "%Y-%m-%d").date()
            if historical_previous_day and (current_day - historical_previous_day).days == 1:
                historical_current += 1
            else:
                historical_current = 1
            historical_max_consecutive = max(historical_max_consecutive, historical_current)
            historical_previous_day = current_day

        reasons = []
        score = 0
        if total > 5:
            score += 45
            reasons.append(f"More than 5 late records in the selected period ({total}).")
        elif total >= 3:
            score += 25
            reasons.append(f"Repeated lateness: {total} records in the selected period.")
        elif total >= 1:
            score += 10
            reasons.append(f"Recorded late {total} time(s) in the selected period.")

        if max_consecutive >= 2:
            score += 30
            reasons.append(f"Consecutive lateness detected for {max_consecutive} day(s) inside the selected period.")
        elif historical_max_consecutive >= 2:
            score += 15
            reasons.append(f"Historical consecutive lateness detected for {historical_max_consecutive} day(s) up to {to_date}.")

        if pattern_for_scoring:
            score += 30
            repeated_text = ', '.join([
                f"{day} ({details['count']} times: {', '.join(details['latest_dates'])})"
                for day, details in pattern_for_scoring.items()
            ])
            reasons.append(f"Repeated weekday pattern detected across weeks: {repeated_text}.")

        if historical_total > 5 and total <= 5:
            score += 15
            reasons.append(f"Historical total lateness is above 5 days up to {to_date} ({historical_total}).")

        if total == 0 and historical_total == 0:
            risk = "Low"
            recommendation = "No action required. Continue normal monitoring."
            reasons.append("No lateness recorded for this student.")
        elif score >= 70:
            risk = "High"
            recommendation = "Immediate follow-up is recommended. Review repeated weekday pattern and assign targeted morning monitoring."
        elif score >= 35:
            risk = "Medium"
            recommendation = "Monitor closely for the next week and check whether the repeated weekday pattern continues."
        else:
            risk = "Low"
            recommendation = "Keep monitoring. No urgent intervention is required now."

        report.append({
            "student_id": st['id'],
            "student_name": st['student_name'],
            "grade_name": st['grade_name'],
            "section_name": st['section_name'],
            "total_late_days": total,
            "historical_total_late_days": historical_total,
            "max_consecutive": max_consecutive,
            "historical_max_consecutive": historical_max_consecutive,
            "repeated_weekdays": repeated_weekdays,
            "historical_repeated_weekdays": historical_repeated_weekdays,
            "risk": risk,
            "score": min(score, 100),
            "reasons": reasons,
            "recommendation": recommendation,
            "records": records,
            "historical_records": historical_records,
        })
    risk_order = {"High": 0, "Medium": 1, "Low": 2}
    report.sort(key=lambda x: (risk_order.get(x['risk'], 3), -x['score'], -x['historical_total_late_days'], -x['total_late_days'], x['grade_name'], x['section_name'], x['student_name']))
    return report


def ai_summary_text(report, from_date, to_date):
    total_students = len(report)
    high = sum(1 for r in report if r['risk'] == 'High')
    medium = sum(1 for r in report if r['risk'] == 'Medium')
    low = sum(1 for r in report if r['risk'] == 'Low')
    total_records = sum(r['total_late_days'] for r in report)
    top = [r for r in report if r['total_late_days'] > 0][:5]
    lines = [
        f"AI Summary for {from_date} to {to_date}",
        f"Total students analyzed: {total_students}.",
        f"Total late records: {total_records}.",
        f"Risk distribution: High {high}, Medium {medium}, Low {low}.",
    ]
    if top:
        lines.append("Top students needing follow-up: " + "; ".join([f"{r['student_name']} ({r['risk']}, {r['total_late_days']} days)" for r in top]) + ".")
    else:
        lines.append("No late records found in this period.")
    return " ".join(lines)


@app.route('/ai-report')
@login_required
def ai_report():
    today_iso = date.today().isoformat()
    from_date = parse_date_or_today(request.args.get('from_date') or today_iso)
    to_date = parse_date_or_today(request.args.get('to_date') or from_date)
    if from_date > to_date:
        from_date, to_date = to_date, from_date
    page = max(request.args.get("page", 1, type=int), 1)
    per_page = request.args.get("per_page", 50, type=int)
    if per_page not in (30, 50, 100):
        per_page = 50
    report_all = get_ai_risk_report(from_date, to_date)
    total_records = len(report_all)
    total_pages = max((total_records + per_page - 1) // per_page, 1)
    if page > total_pages:
        page = total_pages
    start = (page - 1) * per_page
    report = report_all[start:start + per_page]
    summary = ai_summary_text(report_all, from_date, to_date)
    return render_template('ai_report.html', report=report, summary=summary, from_date=from_date, to_date=to_date, today=today_iso, page=page, per_page=per_page, total_pages=total_pages, total_records=total_records, start_index=start)


def is_arabic(text):
    return any('\u0600' <= ch <= '\u06FF' for ch in text or '')


def normalize_text(text):
    text = (text or '').strip().lower()
    text = re.sub(r'[\u064B-\u065F\u0670]', '', text)
    return text


def detect_requested_range(question):
    q = normalize_text(question)
    today = date.today()
    if any(x in q for x in ['today', 'اليوم']):
        return today.isoformat(), today.isoformat()
    if any(x in q for x in ['yesterday', 'امس', 'أمس']):
        d = today - timedelta(days=1)
        return d.isoformat(), d.isoformat()
    if any(x in q for x in ['this week', 'الاسبوع', 'الأسبوع', 'هذا الاسبوع', 'هذا الأسبوع']):
        start = today - timedelta(days=today.weekday())
        return start.isoformat(), today.isoformat()
    if any(x in q for x in ['last week', 'الاسبوع الماضي', 'الأسبوع الماضي']):
        end = today - timedelta(days=today.weekday()+1)
        start = end - timedelta(days=6)
        return start.isoformat(), end.isoformat()
    if any(x in q for x in ['last month', 'previous month', 'الشهر الماضي']):
        first_this = today.replace(day=1)
        last_prev = first_this - timedelta(days=1)
        start_prev = last_prev.replace(day=1)
        return start_prev.isoformat(), last_prev.isoformat()
    if any(x in q for x in ['this month', 'الشهر', 'هذا الشهر']):
        start = today.replace(day=1)
        return start.isoformat(), today.isoformat()
    dates = re.findall(r'\d{4}-\d{2}-\d{2}', question or '')
    if len(dates) >= 2:
        a, b = dates[0], dates[1]
        return (a, b) if a <= b else (b, a)
    if len(dates) == 1:
        return dates[0], dates[0]
    return (today - timedelta(days=30)).isoformat(), today.isoformat()


def get_system_statistics(from_date, to_date):
    report = get_ai_risk_report(from_date, to_date)
    with get_conn() as conn:
        grade_rows = conn.execute('''
            SELECT g.name AS grade_name, COUNT(lr.id) AS total_records,
                   COUNT(DISTINCT lr.student_id) AS unique_students
            FROM grades g
            LEFT JOIN students s ON s.grade_id = g.id AND s.active = 1
            LEFT JOIN late_records lr ON lr.student_id = s.id AND lr.late_date BETWEEN ? AND ?
            GROUP BY g.id
            ORDER BY g.sort_order, g.name
        ''', (from_date, to_date)).fetchall()
        section_rows = conn.execute('''
            SELECT g.name AS grade_name, sec.name AS section_name, COUNT(lr.id) AS total_records,
                   COUNT(DISTINCT lr.student_id) AS unique_students
            FROM sections sec
            JOIN grades g ON g.id = sec.grade_id
            LEFT JOIN students s ON s.section_id = sec.id AND s.active = 1
            LEFT JOIN late_records lr ON lr.student_id = s.id AND lr.late_date BETWEEN ? AND ?
            GROUP BY sec.id
            ORDER BY g.sort_order, sec.name
        ''', (from_date, to_date)).fetchall()
        day_rows = conn.execute('''
            SELECT late_date, COUNT(*) AS total_records
            FROM late_records
            WHERE late_date BETWEEN ? AND ?
            GROUP BY late_date
            ORDER BY late_date
        ''', (from_date, to_date)).fetchall()
        recorder_rows = conn.execute('''
            SELECT COALESCE(u.full_name, u.username, 'Unknown') AS recorder_name, COUNT(lr.id) AS total_records
            FROM late_records lr
            LEFT JOIN users u ON u.id = lr.recorded_by
            WHERE lr.late_date BETWEEN ? AND ?
            GROUP BY recorder_name
            ORDER BY total_records DESC
        ''', (from_date, to_date)).fetchall()
    return {
        'report': report,
        'grades': [dict(x) for x in grade_rows],
        'sections': [dict(x) for x in section_rows],
        'days': [dict(x) for x in day_rows],
        'recorders': [dict(x) for x in recorder_rows],
        'total_records': sum(r['total_late_days'] for r in report),
        'high_risk': sum(1 for r in report if r['risk'] == 'High'),
        'medium_risk': sum(1 for r in report if r['risk'] == 'Medium'),
        'low_risk': sum(1 for r in report if r['risk'] == 'Low'),
    }


def format_today_answer(arabic=False):
    now = datetime.now()
    en_day = calendar.day_name[now.weekday()]
    ar_days = {
        'Monday':'الاثنين', 'Tuesday':'الثلاثاء', 'Wednesday':'الأربعاء',
        'Thursday':'الخميس', 'Friday':'الجمعة', 'Saturday':'السبت', 'Sunday':'الأحد'
    }
    if arabic:
        return f"اليوم هو {ar_days.get(en_day, en_day)}، التاريخ {now.strftime('%Y-%m-%d')}، والوقت الحالي {now.strftime('%H:%M')} حسب وقت السيرفر."
    return f"Today is {en_day}, {now.strftime('%Y-%m-%d')}. Current server time is {now.strftime('%H:%M')}."


def find_student_in_question(question, report):
    q_low = normalize_text(question)
    best = None
    best_score = 0
    for r in report:
        name = normalize_text(r['student_name'])
        if not name:
            continue
        if name in q_low:
            return r
        parts = [p for p in name.split() if len(p) > 2]
        score = sum(1 for p in parts if p in q_low)
        if score > best_score:
            best = r
            best_score = score
    return best if best_score >= 1 else None


def format_student_answer(r, arabic=False):
    if r.get('historical_repeated_weekdays'):
        weekdays = ', '.join([
            f"{d} ({details['count']} times: {', '.join(details.get('latest_dates', []))})"
            for d, details in r['historical_repeated_weekdays'].items()
        ])
    else:
        weekdays = ', '.join([f"{d} ({c})" for d, c in r['repeated_weekdays'].items()]) or ('لا يوجد نمط يوم متكرر' if arabic else 'No repeated weekday pattern')
    recent_records = r.get('records', [])[:10]
    record_lines = []
    for rec in recent_records:
        if arabic:
            record_lines.append(f"- {rec['late_date']} الساعة {rec['late_time']} بواسطة {rec['recorder_name'] or 'غير معروف'}")
        else:
            record_lines.append(f"- {rec['late_date']} at {rec['late_time']} by {rec['recorder_name'] or 'Unknown'}")
    if arabic:
        risk_ar = {'High': 'مرتفع', 'Medium': 'متوسط', 'Low': 'منخفض'}.get(r['risk'], r['risk'])
        return (f"الطالب: {r['student_name']}\n"
                f"الصف/الشعبة: {r['grade_name']} - {r['section_name']}\n"
                f"مستوى الخطورة: {risk_ar} ({r['score']}/100)\n"
                f"مجموع أيام التأخير في الفترة المطلوبة: {r['total_late_days']}\n"
                f"مجموع أيام التأخير التاريخي: {r.get('historical_total_late_days', r['total_late_days'])}\n"
                f"أعلى تأخير متتالي: {r['max_consecutive']} يوم / تاريخيًا {r.get('historical_max_consecutive', r['max_consecutive'])}\n"
                f"نمط الأيام المتكررة بين الأسابيع: {weekdays}\n"
                f"التوصية: {r['recommendation']}\n"
                f"آخر السجلات:\n" + ('\n'.join(record_lines) if record_lines else 'لا توجد سجلات في هذه الفترة.'))
    return (f"Student: {r['student_name']}\n"
            f"Grade/Section: {r['grade_name']} - {r['section_name']}\n"
            f"Risk Level: {r['risk']} ({r['score']}/100)\n"
            f"Total late days in the selected period: {r['total_late_days']}\n"
            f"Historical total late days: {r.get('historical_total_late_days', r['total_late_days'])}\n"
            f"Maximum consecutive late days: {r['max_consecutive']} / historical {r.get('historical_max_consecutive', r['max_consecutive'])}\n"
            f"Repeated weekday pattern across weeks: {weekdays}\n"
            f"Recommendation: {r['recommendation']}\n"
            f"Recent records:\n" + ('\n'.join(record_lines) if record_lines else 'No records in this period.'))


def build_ai_context(question, from_date, to_date):
    stats = get_system_statistics(from_date, to_date)
    report = stats['report']
    mentioned = find_student_in_question(question, report)
    top_students = [r for r in report if r['total_late_days'] > 0][:12]
    high = [r for r in report if r['risk'] == 'High'][:12]
    context_lines = [
        f"Selected period: {from_date} to {to_date}",
        f"Total late records: {stats['total_records']}",
        f"Risk distribution: High={stats['high_risk']}, Medium={stats['medium_risk']}, Low={stats['low_risk']}",
        "Grade totals: " + '; '.join([f"{g['grade_name']}: {g['total_records']} records, {g['unique_students']} students" for g in stats['grades']]),
        "Section totals: " + '; '.join([f"{s['grade_name']} {s['section_name']}: {s['total_records']}" for s in stats['sections'] if s['total_records']]),
        "Daily totals: " + '; '.join([f"{d['late_date']}: {d['total_records']}" for d in stats['days'][:60]]),
        "Recorded by: " + '; '.join([f"{r['recorder_name']}: {r['total_records']}" for r in stats['recorders']]),
        "Top late students: " + '; '.join([f"{r['student_name']} ({r['grade_name']} {r['section_name']}): {r['total_late_days']} days, {r['risk']} risk, consecutive {r['max_consecutive']}" for r in top_students]),
        "High risk students: " + '; '.join([f"{r['student_name']} ({r['total_late_days']} days)" for r in high])
    ]
    if mentioned:
        context_lines.append("Mentioned student details:\n" + format_student_answer(mentioned, arabic=False))
    return '\n'.join(context_lines), stats, mentioned


def openai_smart_answer(question, from_date, to_date):
    api_key = os.environ.get('OPENAI_API_KEY')
    if not api_key or OpenAI is None:
        return None
    context, stats, mentioned = build_ai_context(question, from_date, to_date)
    model = os.environ.get('OPENAI_MODEL', 'gpt-5.5')
    client = OpenAI(api_key=api_key)
    instructions = (
        "You are a smart bilingual Arabic-English AI assistant inside a school lateness monitoring system. "
        "Answer naturally in the same language as the user's question. "
        "Use the provided database context when the question is about students, grades, sections, lateness, risk, users, dates, or reports. "
        "For general questions not related to the database, answer normally and helpfully. "
        "Do not invent student records. If the database context does not contain the requested student or data, say that clearly. "
        "Be concise, professional, and school-administration friendly."
    )
    user_input = f"Database context:\n{context}\n\nUser question:\n{question}"
    response = client.responses.create(model=model, instructions=instructions, input=user_input)
    return response.output_text


def local_smart_answer(question, from_date, to_date):
    q = (question or '').strip()
    q_low = normalize_text(q)
    arabic = is_arabic(q)
    stats = get_system_statistics(from_date, to_date)
    report = stats['report']
    mentioned = find_student_in_question(q, report)

    if any(x in q_low for x in ['what day', 'which day', 'today', 'date', 'time', 'ما هو اليوم', 'ما اليوم', 'اي يوم', 'أي يوم', 'التاريخ', 'الوقت']):
        return format_today_answer(arabic)

    greetings = ['hello', 'hi', 'مرحبا', 'السلام عليكم', 'هلا', 'اهلا', 'أهلا']
    if any(g in q_low for g in greetings) and len(q_low) < 40:
        return 'أهلًا، أنا المساعد الذكي للنظام. اسألني عن الطلاب، التأخير، الخطورة، الصفوف، التقارير، أو التاريخ.' if arabic else 'Hello, I am the smart assistant for this system. Ask me about students, lateness, risk, grades, reports, or dates.'

    if mentioned:
        return format_student_answer(mentioned, arabic)

    # Dynamic section-level smart reports, used by chatbot quick buttons and typed questions.
    with get_conn() as _conn:
        _sections = _conn.execute('''
            SELECT s.name AS section_name, g.name AS grade_name
            FROM sections s JOIN grades g ON g.id = s.grade_id
        ''').fetchall()
    for sec in _sections:
        full_key = normalize_text(f"{sec['grade_name']} {sec['section_name']}")
        sec_key = normalize_text(sec['section_name'])
        if full_key in q_low or (sec_key in q_low and any(w in q_low for w in ['section', 'شعبة', 'الشعبة', 'summary', 'late', 'ملخص', 'تأخير'])):
            sec_rows = [r for r in report if r['grade_name'] == sec['grade_name'] and r['section_name'] == sec['section_name']]
            return format_filtered_report(sec_rows, from_date, to_date, f"{sec['grade_name']} - {sec['section_name']} report", arabic)

    if any(w in q_low for w in ['highest', 'most', 'top', 'late the most', 'أكثر', 'اعلى', 'أعلى']):
        top = [r for r in report if r['total_late_days'] > 0][:10]
        if not top:
            return 'لا توجد سجلات تأخير في الفترة المحددة.' if arabic else 'No late records were found in the selected period.'
        if arabic:
            return 'أكثر الطلاب تأخيرًا في الفترة المحددة:\n' + '\n'.join([f"- {r['student_name']} | {r['grade_name']} {r['section_name']} | {r['total_late_days']} أيام | خطورة {r['risk']}" for r in top])
        return 'Top late students in the selected period:\n' + '\n'.join([f"- {r['student_name']} | {r['grade_name']} {r['section_name']} | {r['total_late_days']} days | {r['risk']} risk" for r in top])

    if any(w in q_low for w in ['risk', 'danger', 'خطورة', 'ريسك', 'خطر', 'high risk', 'عالي الخطورة']):
        filtered_report = report
        grade_label = ''
        grade_map_for_risk = {
            '9': ['grade 9', 'تاسع', 'التاسع', 'صف 9', 'الصف التاسع'],
            '10': ['grade 10', 'عاشر', 'العاشر', 'صف 10', 'الصف العاشر'],
            '11': ['grade 11', 'حادي عشر', 'الحادي عشر', 'صف 11'],
            '12': ['grade 12', 'ثاني عشر', 'الثاني عشر', 'صف 12'],
        }
        for _num, _keys in grade_map_for_risk.items():
            if any(k in q_low for k in _keys):
                filtered_report = [r for r in report if str(grade_number_from_name(r['grade_name'])) == _num]
                grade_label = f" Grade {_num}"
                break
        high = [r for r in filtered_report if r['risk'] == 'High']
        medium = [r for r in filtered_report if r['risk'] == 'Medium']
        if arabic:
            return f"تحليل الخطورة{grade_label} من {from_date} إلى {to_date}: مرتفع {len(high)} طالب، متوسط {len(medium)} طالب.\n" + ('الطلاب عالي الخطورة:\n' + '\n'.join([f"- {r['student_name']} ({r['grade_name']} {r['section_name']} - {r['total_late_days']} أيام)" for r in high[:15]]) if high else 'لا يوجد طلاب عالي الخطورة.')
        return f"Risk analysis{grade_label} from {from_date} to {to_date}: High risk {len(high)} students, Medium risk {len(medium)} students.\n" + ('High-risk students:\n' + '\n'.join([f"- {r['student_name']} ({r['grade_name']} {r['section_name']} - {r['total_late_days']} days)" for r in high[:15]]) if high else 'No high-risk students found.')


    if any(w in q_low for w in ['summary', 'report', 'ملخص', 'تقرير', 'overview']):
        base = ai_summary_text(report, from_date, to_date)
        if arabic:
            return (f"ملخص الفترة من {from_date} إلى {to_date}:\n"
                    f"إجمالي سجلات التأخير: {stats['total_records']}\n"
                    f"عالي الخطورة: {stats['high_risk']}، متوسط: {stats['medium_risk']}، منخفض: {stats['low_risk']}\n"
                    f"حسب الصفوف:\n" + '\n'.join([f"- {g['grade_name']}: {g['total_records']} سجل" for g in stats['grades']]))
        return base + "\nGrade totals:\n" + '\n'.join([f"- {g['grade_name']}: {g['total_records']} records" for g in stats['grades']])

    if any(w in q_low for w in ['who recorded', 'recorded by', 'user', 'users', 'من سجل', 'المستخدم', 'اليوزر']):
        if arabic:
            return 'سجلات الإدخال حسب المستخدم:\n' + ('\n'.join([f"- {r['recorder_name']}: {r['total_records']} سجل" for r in stats['recorders']]) or 'لا توجد سجلات في الفترة المحددة.')
        return 'Records by user:\n' + ('\n'.join([f"- {r['recorder_name']}: {r['total_records']} records" for r in stats['recorders']]) or 'No records in the selected period.')

    grade_map = {
        '9': ['grade 9', 'تاسع', 'التاسع', 'صف 9', 'الصف التاسع'],
        '10': ['grade 10', 'عاشر', 'العاشر', 'صف 10', 'الصف العاشر'],
        '11': ['grade 11', 'حادي عشر', 'الحادي عشر', 'صف 11'],
        '12': ['grade 12', 'ثاني عشر', 'الثاني عشر', 'صف 12'],
    }
    for num, keys in grade_map.items():
        if any(k in q_low for k in keys):
            grade_rows = [r for r in report if str(grade_number_from_name(r['grade_name'])) == num]
            return format_filtered_report(grade_rows, from_date, to_date, f"Grade {num} report", arabic)

    if arabic:
        return ("أستطيع الإجابة عن الأسئلة العامة البسيطة وعن بيانات النظام. للحصول على شات بوت يجيب على أي سؤال عام مثل ChatGPT، "
                "أضف OPENAI_API_KEY في Render Environment. بعد ذلك سأجيب على أي سؤال عام وأحلل بيانات الطلاب بذكاء أعلى.\n\n"
                "جرّب مثلًا: ما هو اليوم؟ من أكثر الطلاب تأخيرًا؟ من سجل التأخير؟ أعطني تقرير الصف العاشر؟")
    return ("I can answer simple general questions and system-data questions locally. To make this chatbot answer any general question like ChatGPT, "
            "add OPENAI_API_KEY in Render Environment. Then it will answer general questions and analyze student data much more intelligently.\n\n"
            "Try: What day is today? Who is late the most? Who recorded the lateness? Give me Grade 10 report.")


def chatbot_answer(question):
    q = (question or '').strip()
    if not q:
        return "Please type a question. / الرجاء كتابة السؤال."
    from_date, to_date = detect_requested_range(q)
    try:
        ai_answer = openai_smart_answer(q, from_date, to_date)
        if ai_answer:
            return ai_answer
    except Exception as exc:
        fallback = local_smart_answer(q, from_date, to_date)
        return fallback + f"\n\n[AI API note: {str(exc)[:180]}]"
    return local_smart_answer(q, from_date, to_date)



def get_chatbot_quick_buttons():
    """Build dynamic clickable prompts for all grades and sections."""
    with get_conn() as conn:
        grades = conn.execute("SELECT id, name FROM grades ORDER BY sort_order, name").fetchall()
        sections = conn.execute('''
            SELECT s.id, s.name AS section_name, g.name AS grade_name
            FROM sections s
            JOIN grades g ON g.id = s.grade_id
            ORDER BY g.sort_order, g.name, s.name
        ''').fetchall()
    grade_buttons = []
    for g in grades:
        grade_name = g['name']
        grade_buttons.extend([
            {"label": f"{grade_name} summary", "question": f"Give me {grade_name} summary this month"},
            {"label": f"{grade_name} late today", "question": f"Show {grade_name} late today"},
            {"label": f"{grade_name} late last week", "question": f"Show {grade_name} late last week"},
            {"label": f"{grade_name} late last month", "question": f"Show {grade_name} late last month"},
            {"label": f"{grade_name} high risk", "question": f"Show high risk students in {grade_name} this month"},
        ])
    section_buttons = []
    for sec in sections:
        full = f"{sec['grade_name']} {sec['section_name']}"
        section_buttons.extend([
            {"label": f"{sec['section_name']} today", "question": f"Show {full} late today"},
            {"label": f"{sec['section_name']} last month", "question": f"Show {full} late last month"},
            {"label": f"{sec['section_name']} summary", "question": f"Give me {full} summary this month"},
        ])
    general_buttons = [
        {"label": "Today late students", "question": "Show late students today"},
        {"label": "Most late students", "question": "Who are the most late students this month?"},
        {"label": "High risk students", "question": "Show high risk students this month"},
        {"label": "Weekly summary", "question": "Give me a summary for this week"},
        {"label": "Last month summary", "question": "Give me a summary for last month"},
        {"label": "Who recorded lateness?", "question": "Who recorded the lateness this month?"},
        {"label": "ما هو اليوم؟", "question": "ما هو اليوم؟"},
        {"label": "ملخص اليوم", "question": "أعطني ملخص التأخير اليوم"},
    ]
    return {"general": general_buttons, "grades": grade_buttons, "sections": section_buttons}


def format_filtered_report(rows, from_date, to_date, title, arabic=False):
    rows = [r for r in rows if r['total_late_days'] > 0]
    total = sum(r['total_late_days'] for r in rows)
    high = sum(1 for r in rows if r['risk'] == 'High')
    top = sorted(rows, key=lambda r: (-r['total_late_days'], r['student_name']))[:15]
    if arabic:
        if not top:
            return f"{title}\nالفترة: {from_date} إلى {to_date}\nلا توجد سجلات تأخير في هذه الفترة."
        return (f"{title}\nالفترة: {from_date} إلى {to_date}\n"
                f"إجمالي سجلات التأخير: {total}\n"
                f"عدد الطلاب المتأخرين: {len(rows)}\n"
                f"طلاب عالي الخطورة: {high}\n\n"
                "أكثر الطلاب تأخيرًا:\n" +
                "\n".join([f"- {r['student_name']} | {r['grade_name']} {r['section_name']} | {r['total_late_days']} أيام | {r['risk']}" for r in top]))
    if not top:
        return f"{title}\nPeriod: {from_date} to {to_date}\nNo late records were found in this period."
    return (f"{title}\nPeriod: {from_date} to {to_date}\n"
            f"Total late records: {total}\n"
            f"Late students: {len(rows)}\n"
            f"High-risk students: {high}\n\n"
            "Most late students:\n" +
            "\n".join([f"- {r['student_name']} | {r['grade_name']} {r['section_name']} | {r['total_late_days']} days | {r['risk']}" for r in top]))


@app.route('/chatbot', methods=['GET', 'POST'])
@login_required
def chatbot():
    answer = None
    question = ''
    quick_buttons = get_chatbot_quick_buttons()
    if request.method == 'POST':
        question = request.form.get('quick_question') or request.form.get('question', '')
        answer = chatbot_answer(question)
    return render_template('chatbot.html', question=question, answer=answer, quick_buttons=quick_buttons)


def get_smtp_settings():
    smtp_host = os.environ.get("SMTP_HOST", "smtp.gmail.com")
    smtp_port = int(os.environ.get("SMTP_PORT", "587"))
    smtp_user = os.environ.get("SMTP_USER") or os.environ.get("MAIL_USERNAME")
    smtp_password = os.environ.get("SMTP_PASSWORD") or os.environ.get("MAIL_PASSWORD")
    sender = os.environ.get("SMTP_FROM") or os.environ.get("EMAIL_FROM") or smtp_user or ""
    return smtp_host, smtp_port, smtp_user, smtp_password, sender


def build_email_body(day_text, grade_name=None):
    title_grade = f" - {grade_name}" if grade_name else ""
    records = get_records_range(day_text, day_text)
    if grade_name:
        records = [r for r in records if r["grade_name"] == grade_name]
    if not records:
        return (
            f"Daily Late/Absent Report{title_grade} - {day_text}\n\n"
            "No late/absent students recorded today.\n\n"
            "This email was generated automatically by the Smart Late Students System."
        )
    lines = [
        f"Daily Late/Absent Report{title_grade} - {day_text}",
        "",
        "Student | Grade | Section | Date | Time | Total Days | Recorded By",
        "-" * 95,
    ]
    for r in records:
        lines.append(
            f"{r['student_name']} | {r['grade_name']} | {r['section_name']} | "
            f"{r['late_date']} | {r['late_time']} | {r['total_days']} | "
            f"{r['recorder_name'] or 'Unknown'}"
        )
    lines.append("\nAttached: official Excel report for printing/submission.")
    return "\n".join(lines)

def _attachment_b64(attachment_bytes):
    """Return base64 string for API email attachments."""
    raw = attachment_bytes.getvalue() if hasattr(attachment_bytes, "getvalue") else attachment_bytes
    return base64.b64encode(raw).decode("utf-8")


def get_email_provider():
    """Choose the safest email provider for Render.

    Preferred provider for this project:
      EMAIL_PROVIDER=brevo + BREVO_API_KEY + EMAIL_FROM

    Resend, SendGrid, and SMTP remain as optional fallbacks only.
    SMTP is not recommended on Render because some Render networks cannot reach
    smtp.gmail.com ports 465/587.
    """
    provider = (os.environ.get("EMAIL_PROVIDER") or "").strip().lower()
    if provider:
        return provider
    if os.environ.get("BREVO_API_KEY"):
        return "brevo"
    if os.environ.get("RESEND_API_KEY"):
        return "resend"
    if os.environ.get("SENDGRID_API_KEY"):
        return "sendgrid"
    return "smtp"


def _parse_sender(sender_text):
    """Return (name, email) from values like 'Attendance System <mail@example.com>'."""
    sender_text = (sender_text or "").strip()
    if "<" in sender_text and ">" in sender_text:
        name = sender_text.split("<", 1)[0].strip().strip('"') or "Attendance System"
        email = sender_text.split("<", 1)[1].split(">", 1)[0].strip()
        return name, email
    return "Attendance System", sender_text


def send_email_brevo(to_email, subject, body, attachment_bytes, filename):
    """Send email using Brevo Transactional Email API. Recommended on Render."""
    api_key = os.environ.get("BREVO_API_KEY") or os.environ.get("SENDINBLUE_API_KEY")
    sender = os.environ.get("EMAIL_FROM") or os.environ.get("BREVO_FROM") or os.environ.get("SMTP_FROM")
    if not api_key or not sender:
        return False, "Brevo settings are missing. Add BREVO_API_KEY and EMAIL_FROM in Render Environment."

    sender_name, sender_email = _parse_sender(sender)
    if not sender_email or "@" not in sender_email:
        return False, "EMAIL_FROM is invalid. Use format: Attendance System <your_verified_sender@email.com>"

    payload = {
        "sender": {"name": sender_name, "email": sender_email},
        "to": [{"email": to_email}],
        "subject": subject,
        "textContent": body,
        "attachment": [
            {
                "name": filename,
                "content": _attachment_b64(attachment_bytes),
            }
        ],
    }

    try:
        response = requests.post(
            "https://api.brevo.com/v3/smtp/email",
            headers={
                "api-key": api_key,
                "accept": "application/json",
                "Content-Type": "application/json",
            },
            json=payload,
            timeout=30,
        )
        if 200 <= response.status_code < 300:
            return True, "Email sent successfully using Brevo API."
        return False, f"Brevo API error {response.status_code}: {response.text[:400]}"
    except requests.RequestException as exc:
        return False, f"Brevo connection failed: {str(exc)[:220]}"


def send_email_resend(to_email, subject, body, attachment_bytes, filename):
    api_key = os.environ.get("RESEND_API_KEY")
    sender = os.environ.get("EMAIL_FROM") or os.environ.get("SMTP_FROM") or os.environ.get("RESEND_FROM")
    if not api_key or not sender:
        return False, "Resend settings are missing. Add RESEND_API_KEY and EMAIL_FROM in Render Environment."

    payload = {
        "from": sender,
        "to": [to_email],
        "subject": subject,
        "text": body,
        "attachments": [
            {
                "filename": filename,
                "content": _attachment_b64(attachment_bytes),
            }
        ],
    }
    try:
        response = requests.post(
            "https://api.resend.com/emails",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json=payload,
            timeout=25,
        )
        if 200 <= response.status_code < 300:
            return True, "Email sent successfully using Resend API."
        return False, f"Resend API error {response.status_code}: {response.text[:300]}"
    except requests.RequestException as exc:
        return False, f"Resend connection failed: {str(exc)[:220]}"


def send_email_sendgrid(to_email, subject, body, attachment_bytes, filename):
    api_key = os.environ.get("SENDGRID_API_KEY")
    sender = os.environ.get("EMAIL_FROM") or os.environ.get("SMTP_FROM") or os.environ.get("SENDGRID_FROM")
    if not api_key or not sender:
        return False, "SendGrid settings are missing. Add SENDGRID_API_KEY and EMAIL_FROM in Render Environment."

    payload = {
        "personalizations": [{"to": [{"email": to_email}]}],
        "from": {"email": sender},
        "subject": subject,
        "content": [{"type": "text/plain", "value": body}],
        "attachments": [
            {
                "content": _attachment_b64(attachment_bytes),
                "type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "filename": filename,
                "disposition": "attachment",
            }
        ],
    }
    try:
        response = requests.post(
            "https://api.sendgrid.com/v3/mail/send",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json=payload,
            timeout=25,
        )
        if 200 <= response.status_code < 300:
            return True, "Email sent successfully using SendGrid API."
        return False, f"SendGrid API error {response.status_code}: {response.text[:300]}"
    except requests.RequestException as exc:
        return False, f"SendGrid connection failed: {str(exc)[:220]}"


def send_email_smtp(to_email, subject, body, attachment_bytes, filename):
    """Fallback SMTP sender. API providers are recommended on Render."""
    smtp_host, smtp_port, smtp_user, smtp_password, sender = get_smtp_settings()
    if not smtp_user or not smtp_password or not sender:
        return False, "SMTP settings are missing. Add SMTP_USER, SMTP_PASSWORD, SMTP_HOST, SMTP_PORT and SMTP_FROM in Render Environment, or use BREVO_API_KEY."

    msg = EmailMessage()
    msg["Subject"] = subject
    msg["From"] = sender
    msg["To"] = to_email
    msg.set_content(body)
    msg.add_attachment(
        attachment_bytes.getvalue(),
        maintype="application",
        subtype="vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        filename=filename,
    )

    try:
        if int(smtp_port) == 465:
            with smtplib.SMTP_SSL(smtp_host, int(smtp_port), timeout=15) as server:
                server.login(smtp_user, smtp_password)
                server.send_message(msg)
        else:
            with smtplib.SMTP(smtp_host, int(smtp_port), timeout=15) as server:
                server.ehlo()
                server.starttls()
                server.ehlo()
                server.login(smtp_user, smtp_password)
                server.send_message(msg)
        return True, "Email sent successfully using SMTP."
    except smtplib.SMTPAuthenticationError:
        return False, "SMTP authentication failed. Check SMTP_USER and Gmail App Password. Do not use your normal Gmail password."
    except (TimeoutError, OSError) as exc:
        return False, f"SMTP connection failed or timed out. Render may block SMTP. Use Brevo API. Details: {str(exc)[:160]}"
    except Exception as exc:
        return False, f"SMTP email sending failed: {type(exc).__name__}: {str(exc)[:180]}"


def send_email_with_attachment(to_email, subject, body, attachment_bytes, filename):
    """Send Excel report using an Email API first, SMTP only as fallback."""
    provider = get_email_provider()
    if provider == "brevo":
        return send_email_brevo(to_email, subject, body, attachment_bytes, filename)
    if provider == "resend":
        return send_email_resend(to_email, subject, body, attachment_bytes, filename)
    if provider == "sendgrid":
        return send_email_sendgrid(to_email, subject, body, attachment_bytes, filename)
    if provider == "smtp":
        return send_email_smtp(to_email, subject, body, attachment_bytes, filename)
    return False, "Invalid EMAIL_PROVIDER. Use brevo, resend, sendgrid, or smtp."


def send_grade_email(grade_id, recipient_id, day_text=None):
    day_text = day_text or date.today().isoformat()
    with get_conn() as conn:
        grade = conn.execute("SELECT * FROM grades WHERE id=?", (int(grade_id),)).fetchone() if grade_id else None
        recipient = conn.execute("SELECT * FROM email_recipients WHERE id=? AND active=1", (int(recipient_id),)).fetchone() if recipient_id else None
    if not grade:
        return False, "Please select a valid grade."
    if not recipient:
        return False, "Please select an active recipient."
    if recipient["grade_id"] != grade["id"]:
        return False, "The selected recipient is not assigned to this grade."
    records = get_records_range(day_text, day_text, grade_id=grade["id"])
    excel_bytes = build_excel_workbook_bytes(records, day_text, day_text, grade_id=grade["id"])
    safe_grade = sanitize_sheet_name(grade["name"]).replace(" ", "_")
    filename = f"late_absent_{safe_grade}_{day_text}.xlsx"
    subject = f"Daily Late/Absent Report - {grade['name']} - {day_text}"
    body = build_email_body(day_text, grade_name=grade["name"])
    return send_email_with_attachment(recipient["email"], subject, body, excel_bytes, filename)


def send_daily_email():
    """Send one daily Excel email per active grade recipient.

    Schedule this through Render Cron Job at 08:30 using: python cron.py
    """
    day_text = date.today().isoformat()
    sent = 0
    errors = []
    with get_conn() as conn:
        recipients = conn.execute("""
            SELECT er.*, g.name AS grade_name
            FROM email_recipients er JOIN grades g ON g.id=er.grade_id
            WHERE er.active=1
            ORDER BY g.sort_order, er.person_name
        """).fetchall()
    if not recipients:
        print("No active email recipients configured. Email skipped.")
        return
    for rec in recipients:
        ok, message = send_grade_email(rec["grade_id"], rec["id"], day_text)
        if ok:
            sent += 1
            print(f"Sent daily report to {rec['person_name']} <{rec['email']}> for {rec['grade_name']}")
        else:
            errors.append(f"{rec['email']}: {message}")
            print(f"Failed daily report to {rec['email']}: {message}")
    print(f"Daily email complete. Sent: {sent}. Errors: {len(errors)}")



# =========================
# ADMIN ANALYTICS DASHBOARD + PPT EXPORT
# =========================

def iso_to_date(value):
    return datetime.strptime(value, "%Y-%m-%d").date()


def safe_text(value, limit=180):
    text = str(value or "-")
    return text if len(text) <= limit else text[:limit-3] + "..."


def group_period_label(day_text, from_date, to_date):
    """Return weekly labels for short/medium intervals and monthly labels for long intervals."""
    d = iso_to_date(day_text)
    start = iso_to_date(from_date)
    end = iso_to_date(to_date)
    if (end - start).days > 62:
        return d.strftime("%Y-%m")
    week_index = ((d - start).days // 7) + 1
    chunk_start = start + timedelta(days=(week_index - 1) * 7)
    chunk_end = min(chunk_start + timedelta(days=6), end)
    return f"Week {week_index}: {chunk_start.isoformat()} to {chunk_end.isoformat()}"


def empty_chart_payload(title):
    return {"title": title, "labels": ["No Data"], "values": [0], "details": [], "tooltip_lines": ["No data available"], "group_summaries": [], "student_summary": []}


def get_admin_analytics(from_date, to_date, grade_id=None, section_id=None):
    records = get_records_range(from_date, to_date, grade_id=grade_id, section_id=section_id)
    period_label = f"{from_date} to {to_date}"
    if not records:
        return {
            "from_date": from_date,
            "to_date": to_date,
            "grade_id": grade_id,
            "section_id": section_id,
            "total_records": 0,
            "unique_students": 0,
            "trend": empty_chart_payload("Weekly / Monthly Trend"),
            "sections": empty_chart_payload("Most Late/Absent Sections"),
            "students": empty_chart_payload("Most Late Students"),
            "recent_records": [],
            "smart_summary": f"No late/absent records found for {period_label}.",
        }

    # 1) Weekly/monthly trend: chart shows late/absent record counts.
    trend = {}
    trend_details = {}
    for r in records:
        label = group_period_label(r["late_date"], from_date, to_date)
        trend[label] = trend.get(label, 0) + 1
        trend_details.setdefault(label, []).append(r)

    # 2) Sections: exact division like 10 ADV 1 / 12 GEN 2.
    sections = {}
    section_details = {}
    for r in records:
        label = f"{r['grade_name']} - {r['section_name']}"
        sections[label] = sections.get(label, 0) + 1
        section_details.setdefault(label, []).append(r)

    # 3) Most late students.
    student_counts = {}
    student_details = {}
    student_display = {}
    for r in records:
        sid = r["student_id"]
        student_counts[sid] = student_counts.get(sid, 0) + 1
        student_details.setdefault(sid, []).append(r)
        student_display[sid] = {
            "student_number": r.get("student_number") or "-",
            "student_name": r.get("student_name_en") or r.get("student_name") or "-",
            "student_name_ar": r.get("student_name_ar") or "-",
            "grade": r.get("grade_name") or "-",
            "section": r.get("section_name") or "-",
        }

    def detail_from_record(group, r):
        return {
            "group": group,
            "student_id": r.get("student_id"),
            "student_number": r.get("student_number") or "-",
            "student_name": r.get("student_name_en") or r.get("student_name") or "-",
            "student_name_ar": r.get("student_name_ar") or "-",
            "grade": r.get("grade_name") or "-",
            "section": r.get("section_name") or "-",
            "date": r.get("late_date") or "-",
            "day": r.get("day_name") or "-",
            "time": r.get("late_time") or "-",
            "recorded_by": r.get("recorder_name") or r.get("recorder_username") or "Unknown",
        }

    def summarize_students_for_details(recs):
        summary = {}
        for r in recs:
            sid = r.get("student_id")
            key = sid or f"{r.get('student_number')}-{r.get('student_name_en')}"
            if key not in summary:
                summary[key] = {
                    "student_id": sid,
                    "student_number": r.get("student_number") or "-",
                    "student_name": r.get("student_name_en") or r.get("student_name") or "-",
                    "student_name_ar": r.get("student_name_ar") or "-",
                    "grade": r.get("grade_name") or "-",
                    "section": r.get("section_name") or "-",
                    "count": 0,
                    "dates": [],
                }
            summary[key]["count"] += 1
            summary[key]["dates"].append(f"{r.get('late_date')} {r.get('late_time') or ''}".strip())
        return sorted(summary.values(), key=lambda x: (-x["count"], x["grade"], x["section"], x["student_name"]))

    def build_group_payload(title, counts, details_map, sort_desc=False, limit=None):
        items = list(counts.items())
        if sort_desc:
            items.sort(key=lambda x: (-x[1], x[0]))
        if limit:
            items = items[:limit]
        labels = [str(k) for k, _ in items] or ["No Data"]
        values = [v for _, v in items] or [0]
        details = []
        tooltip_lines = []
        group_summaries = []
        for label, count in items:
            recs = details_map.get(label, [])
            student_summary = summarize_students_for_details(recs)
            names = [x["student_name"] for x in student_summary[:6]]
            tooltip_lines.append(f"{count} records | " + ", ".join(names) + ("..." if len(student_summary) > 6 else ""))
            group_summaries.append({"label": str(label), "count": count, "students": student_summary})
            for r in recs:
                details.append(detail_from_record(str(label), r))
        return {"title": title, "labels": labels, "values": values, "details": details, "tooltip_lines": tooltip_lines or ["No data available"], "group_summaries": group_summaries}

    def build_students_payload():
        items = sorted(student_counts.items(), key=lambda x: (-x[1], student_display[x[0]]["grade"], student_display[x[0]]["section"], student_display[x[0]]["student_name"]))[:15]
        labels = []
        values = []
        details = []
        tooltip_lines = []
        student_summary = []
        for idx, (sid, count) in enumerate(items, start=1):
            display = student_display[sid]
            readable_label = f"{display['student_name']} ({display['grade']} {display['section']})"
            labels.append(safe_text(readable_label, 34))
            values.append(count)
            records_for_student = sorted(student_details.get(sid, []), key=lambda r: (r["late_date"], r.get("late_time") or ""))
            dates = [f"{r['late_date']} {r.get('late_time','')}".strip() for r in records_for_student]
            student_summary.append({
                "rank": idx,
                "student_id": sid,
                "student_number": display["student_number"],
                "student_name": display["student_name"],
                "student_name_ar": display["student_name_ar"],
                "grade": display["grade"],
                "section": display["section"],
                "total_late_days": count,
                "dates": dates,
            })
            tooltip_lines.append(f"{display['student_name']} | {display['grade']} {display['section']} | {count} records")
            for r in records_for_student:
                details.append(detail_from_record(display['student_name'], r))
        return {"title": "Most Late Students", "labels": labels or ["No Data"], "values": values or [0], "details": details, "tooltip_lines": tooltip_lines or ["No data available"], "student_summary": student_summary}

    top_section = sorted(sections.items(), key=lambda x: x[1], reverse=True)[0] if sections else None
    top_student_sid = sorted(student_counts.items(), key=lambda x: x[1], reverse=True)[0][0] if student_counts else None
    top_student = student_display.get(top_student_sid, {}) if top_student_sid else {}
    smart_summary = f"{len(records)} records found from {from_date} to {to_date}."
    if top_section:
        smart_summary += f" Highest section: {top_section[0]} with {top_section[1]} records."
    if top_student_sid:
        smart_summary += f" Most late student: {top_student.get('student_name', '-')}, {top_student.get('grade', '-')} {top_student.get('section', '-')}, with {student_counts[top_student_sid]} records."

    return {
        "from_date": from_date,
        "to_date": to_date,
        "grade_id": grade_id,
        "section_id": section_id,
        "total_records": len(records),
        "unique_students": len({r["student_id"] for r in records}),
        "trend": build_group_payload("Weekly / Monthly Trend", trend, trend_details),
        "sections": build_group_payload("Most Late/Absent Sections", sections, section_details, sort_desc=True, limit=15),
        "students": build_students_payload(),
        "recent_records": [detail_from_record("Record", r) for r in records[-20:]],
        "smart_summary": smart_summary,
    }

@app.route('/admin/analytics')
@login_required
@admin_required
def admin_analytics():
    today_iso = date.today().isoformat()
    default_from = date.today().replace(day=1).isoformat()
    from_date = parse_date_or_today(request.args.get('from_date') or default_from)
    to_date = parse_date_or_today(request.args.get('to_date') or today_iso)
    if from_date > to_date:
        from_date, to_date = to_date, from_date
    grade_id = request.args.get('grade_id', type=int)
    section_id = request.args.get('section_id', type=int)
    with get_conn() as conn:
        grades = conn.execute("SELECT * FROM grades ORDER BY sort_order, name").fetchall()
        sections = conn.execute("SELECT s.*, g.name AS grade_name FROM sections s JOIN grades g ON g.id=s.grade_id ORDER BY g.sort_order, s.name").fetchall()
    analytics = get_admin_analytics(from_date, to_date, grade_id=grade_id, section_id=section_id)
    return render_template('admin_analytics.html', analytics=analytics, from_date=from_date, to_date=to_date, today=today_iso, grades=grades, sections=sections, selected_grade=grade_id, selected_section=section_id)


def add_title(slide, title, subtitle=None):
    slide.shapes.title.text = safe_text(title, 100)
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = safe_text(subtitle, 160)


def add_ppt_chart_slide(prs, payload, from_date, to_date):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = f"{payload['title']} ({from_date} to {to_date})"
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.size = Pt(22)

    chart_data = CategoryChartData()
    chart_data.categories = [safe_text(x, 40) for x in payload['labels']]
    chart_data.add_series('Late/Absent Records', payload['values'])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.55), Inches(0.95), Inches(12.1), Inches(3.6),
        chart_data
    ).chart
    chart.has_legend = True
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = payload['title']

    # Short detail table under chart. Full detail slides are added separately.
    rows = min(len(payload['details']), 8) + 1
    table_shape = slide.shapes.add_table(rows, 6, Inches(0.55), Inches(4.78), Inches(12.1), Inches(1.65))
    table = table_shape.table
    headers = ["Group", "Student No.", "Student", "Grade", "Section", "Date"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    for row_idx, item in enumerate(payload['details'][:8], start=1):
        values = [item['group'], item['student_number'], item['student_name'], item['grade'], item['section'], item['date']]
        for col_idx, value in enumerate(values):
            table.cell(row_idx, col_idx).text = safe_text(value, 40)


def add_detail_slides(prs, payload, from_date, to_date):
    details = payload.get('details') or []
    if not details:
        return
    headers = ["Group", "Student No.", "English Name", "Arabic Name", "Grade", "Section", "Date", "Day", "Time", "Recorded By"]
    chunk_size = 12
    for idx in range(0, len(details), chunk_size):
        chunk = details[idx:idx + chunk_size]
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.2), Inches(12.7), Inches(0.45))
        title_box.text_frame.text = f"Details - {payload['title']} ({from_date} to {to_date})"
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.size = Pt(18)
        table_shape = slide.shapes.add_table(len(chunk) + 1, len(headers), Inches(0.25), Inches(0.85), Inches(12.85), Inches(6.25))
        table = table_shape.table
        for c, h in enumerate(headers):
            table.cell(0, c).text = h
        for r_idx, item in enumerate(chunk, start=1):
            values = [item['group'], item['student_number'], item['student_name'], item['student_name_ar'], item['grade'], item['section'], item['date'], item['day'], item['time'], item['recorded_by']]
            for c_idx, value in enumerate(values):
                table.cell(r_idx, c_idx).text = safe_text(value, 28)


def build_analytics_pptx(analytics, chart_key):
    if Presentation is None:
        raise RuntimeError("python-pptx is missing. Add python-pptx to requirements.txt.")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Attendance Analytics Dashboard"
    title_slide.placeholders[1].text = f"Period: {analytics['from_date']} to {analytics['to_date']}\nTotal records: {analytics['total_records']} | Students with records: {analytics['unique_students']}"

    chart_keys = [chart_key] if chart_key != "all" else ["trend", "sections", "students"]
    chart_names = {"trend": "Weekly / Monthly Trend", "sections": "Top Sections", "students": "Most Late Students"}
    for key in chart_keys:
        if key not in analytics:
            continue
        payload = analytics[key]
        add_ppt_chart_slide(prs, payload, analytics['from_date'], analytics['to_date'])
        add_detail_slides(prs, payload, analytics['from_date'], analytics['to_date'])

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


@app.route('/admin/analytics/export/<chart_key>.pptx')
@login_required
@admin_required
def export_analytics_pptx(chart_key):
    today_iso = date.today().isoformat()
    default_from = date.today().replace(day=1).isoformat()
    from_date = parse_date_or_today(request.args.get('from_date') or default_from)
    to_date = parse_date_or_today(request.args.get('to_date') or today_iso)
    if from_date > to_date:
        from_date, to_date = to_date, from_date
    if chart_key not in ("trend", "sections", "students", "all"):
        chart_key = "all"
    grade_id = request.args.get("grade_id", type=int)
    section_id = request.args.get("section_id", type=int)
    analytics = get_admin_analytics(from_date, to_date, grade_id=grade_id, section_id=section_id)
    pptx = build_analytics_pptx(analytics, chart_key)
    suffix = from_date if from_date == to_date else f"{from_date}_to_{to_date}"
    return send_file(
        pptx,
        as_attachment=True,
        download_name=f"attendance_analytics_{chart_key}_{suffix}.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@app.route("/send-report-now")
def send_report_now():
    if request.args.get("token") != os.environ.get("CRON_SECRET", "change-me"):
        return Response("Unauthorized", status=401)
    send_daily_email()
    return "Report email sent or skipped. Check logs."


if __name__ == "__main__":
    init_db()
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)), debug=True)
else:
    init_db()
